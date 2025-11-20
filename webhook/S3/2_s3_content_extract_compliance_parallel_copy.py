#**Python Script 2:** Fault-tolerant, parallel content fetch + LLM (OpenAI) classification, resumes on failure.


import os
import sys
import json
import logging
import io
from concurrent.futures import ThreadPoolExecutor, as_completed
from docx import Document
from PyPDF2 import PdfReader
import boto3
import re

# --- NEW IMPORTS FOR EXTRA FILETYPES -- #
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None
try:
    from odf.opendocument import load as odf_load
    from odf.text import P
except ImportError:
    odf_load = None
try:
    import extract_msg
except ImportError:
    extract_msg = None
import email
import email.policy
# --------------------------------------- #

try:
    import openai
except ImportError:
    openai = None

logging.basicConfig(format="%(asctime)s %(levelname)s %(message)s", level=logging.INFO)
ALLOWED_EXT = (
    '.docx', '.txt', '.pdf', '.csv', '.json', '.md',     # base + markdown
    '.xls', '.xlsx', '.pptx',                            # Excel, PPTX
    '.rtf', '.odt', '.ods',                             # Rich Text, OpenDoc
    '.eml', '.msg'                                      # Email file formats
)
MAX_OUT_CHARS = 15000
MAX_WORKERS = 8

def abort(msg, code=1):
    logging.error(msg)
    sys.exit(code)

if len(sys.argv) < 2:
    abort("Usage: python3 2_s3_content_extract_compliance_parallel.py <CONFIG_ID> [<search_root>]")

config_id = sys.argv[1]
search_root = sys.argv[2] if len(sys.argv) > 2 else '/home/cybersecai/htdocs/www.cybersecai.io/webhook/S3'
def find_config_folder(config_id, root):
    for rootdir, dirs, files in os.walk(root):
        if os.path.basename(rootdir) == config_id:
            return os.path.abspath(rootdir)
    abort(f"Could not find config folder '{config_id}' in '{root}'")

def load_json(path, desc):
    if not os.path.isfile(path):
        abort(f"{desc} not found: {path}")
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

config_dir = find_config_folder(config_id, search_root)
graph_dir = os.path.join(config_dir, "graph")
os.makedirs(graph_dir, exist_ok=True)

secrets_json = os.path.join(config_dir, f"{config_id}_secrets.json")
compliance_matrix = load_json(os.path.join(config_dir, "compliance_matrix.json"), "Compliance matrix")
s3config = load_json(secrets_json, "S3 Secrets file")
AWS_ACCESS_KEY = s3config['AWS_ACCESS_KEY']
AWS_SECRET_KEY = s3config['AWS_SECRET_KEY']
region = s3config.get('AWS_REGION') or s3config.get('AWS_REGION', 'us-east-1')  # support older secrets
s3 = boto3.client(
    "s3",
    aws_access_key_id=AWS_ACCESS_KEY,
    aws_secret_access_key=AWS_SECRET_KEY,
    region_name=region
)

def fetch_file_content(bucket, record):
    key = record['key']
    ext = os.path.splitext(record['file_name'])[1].lower()
    if ext not in ALLOWED_EXT or record.get('file_type','') == 'folder':
        return None
    try:
        response = s3.get_object(Bucket=bucket, Key=key)
        if ext in ('.txt', '.csv', '.json', '.md'):
            content = response['Body'].read(MAX_OUT_CHARS*2)
            try:
                return content.decode('utf-8')
            except Exception:
                return content.decode('latin-1', errors="replace")

        elif ext == '.docx':
            tmpname = "/tmp/tmps3_%s.docx" % (os.urandom(4).hex())
            with open(tmpname, "wb") as tmpf:
                tmpf.write(response['Body'].read())
            try:
                doc = Document(tmpname)
                return "\n".join(p.text for p in doc.paragraphs)
            finally:
                try: os.unlink(tmpname)
                except Exception: pass

        elif ext == '.pdf':
            content = response['Body'].read(3_000_000)
            pdf_bytes = io.BytesIO(content)
            reader = PdfReader(pdf_bytes)
            text = ""
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
            return text[:MAX_OUT_CHARS]


        elif ext in ('.xls', '.xlsx'):
            if pd is None:
                logging.warning("Pandas not installed; cannot parse Excel files")
                return None

            tmpname = f"/tmp/tmps3_{os.urandom(4).hex()}{ext}"
            with open(tmpname, "wb") as tmpf:
                tmpf.write(response['Body'].read())

            try:
                engine = 'openpyxl' if ext == '.xlsx' else 'xlrd'
                xls = pd.ExcelFile(tmpname, engine=engine)
                pieces = []
                for sheet in xls.sheet_names:
                    try:
                        df = pd.read_excel(
                            xls, sheet_name=sheet, dtype=str, nrows=1000, na_filter=False
                        )
                        if not df.empty:
                            csv_text = df.to_csv(index=False)
                            if csv_text.strip():
                                pieces.append(f"=== Sheet: {sheet} ===\n{csv_text}")
                    except Exception as se:
                        logging.warning(f"Excel sheet parse error ({sheet}): {se}")
                combined = "\n\n".join(pieces)
                return combined[:MAX_OUT_CHARS] if combined else None
            except Exception as e:
                logging.warning(f"Excel parse error: {e}")
                return None
            finally:
                try:
                    os.unlink(tmpname)
                except Exception:
                    pass

        

        elif ext == '.pptx':
            if Presentation is None:
                logging.warning("python-pptx not installed; cannot parse .pptx")
                return None

            tmpname = f"/tmp/tmps3_{os.urandom(4).hex()}.pptx"
            with open(tmpname, "wb") as tmpf:
                tmpf.write(response['Body'].read())

            try:
                prs = Presentation(tmpname)
                lines = []

                for slide in prs.slides:
                    # slide title
                    if slide.shapes.title and slide.shapes.title.text:
                        lines.append(slide.shapes.title.text)

                    for shape in slide.shapes:
                        # text frames
                        if getattr(shape, "has_text_frame", False) and shape.text:
                            lines.append(shape.text)

                    # tables
                    if hasattr(shape, "table") and shape.table is not None:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    lines.append(cell.text)

                return "\n".join(lines)

            except Exception as e:
                logging.warning(f"PPTX parse error: {e}")
                return None

            finally:
                try:
                    os.unlink(tmpname)
                except Exception:
                    pass
        # elif ext == '.pptx' and Presentation is not None:
        #     #tmpname = "/tmp/tmps3_%s.pptx" % (os.urandom(4).hex())
        #     tmpname = f"/tmp/tmps3_{os.urandom(4).hex()}.pptx"
        #     with open(tmpname, "wb") as tmpf:
        #         tmpf.write(response['Body'].read())
        #     try:
        #         prs = Presentation(tmpname)
        #         lines = []
        #         for slide in prs.slides:
        #             for shape in slide.shapes:
        #                 if hasattr(shape, "text"):
        #                     lines.append(shape.text)
        #         return "\n".join(lines)
        #     except Exception as e:
        #         logging.warning(f"PPTX parse error: {e}")
        #         return None
        #     finally:
        #         try: os.unlink(tmpname)
        #         except Exception: pass

        elif ext == '.rtf' and rtf_to_text is not None:
            content = response['Body'].read()
            try:
                return rtf_to_text(content.decode(errors="ignore"))
            except Exception as e:
                logging.warning(f"RTF parse error: {e}")
                return None

        elif ext in ('.odt', '.ods') and odf_load is not None:
            tmpname = f"/tmp/tmps3_{os.urandom(4).hex()}{ext}"
            with open(tmpname, "wb") as tmpf:
                tmpf.write(response['Body'].read())
            try:
                doc = odf_load(tmpname)
                paragraphs = []
                for elem in doc.getElementsByType(P):
                    paragraphs.append(str(elem))
                return "\n".join(paragraphs)
            except Exception as e:
                logging.warning(f"ODF parse error: {e}")
                return None
            finally:
                try: os.unlink(tmpname)
                except Exception: pass

        elif ext == '.eml':
            content = response['Body'].read(2_000_000)
            try:
                msg = email.message_from_bytes(content, policy=email.policy.default)
                text_parts = []
                if msg['Subject']:
                    text_parts.append("Subject: " + str(msg['Subject']))
                if msg['From']:
                    text_parts.append("From: " + str(msg['From']))
                if msg['To']:
                    text_parts.append("To: " + str(msg['To']))
                if msg['Date']:
                    text_parts.append("Date: " + str(msg['Date']))
                body = ""
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            body += str(part.get_content())
                else:
                    body = msg.get_content()
                text_parts.append(str(body))
                return "\n".join(text_parts)
            except Exception as e:
                logging.warning(f"EML parse error: {e}")
                return None

        elif ext == '.msg':
            if extract_msg is None:
                logging.warning("extract-msg not installed; cannot parse .msg")
                return None

            tmpname = f"/tmp/tmps3_{os.urandom(4).hex()}.msg"
            with open(tmpname, "wb") as tmpf:
                tmpf.write(response['Body'].read())

            try:
                msg = extract_msg.Message(tmpname)
                subject = msg.subject or ""
                sender = getattr(msg, "sender", "") or getattr(msg, "sender_email", "")
                to = msg.to or ""
                date = msg.date or ""
                body = msg.body or ""
        
                # Fallback to HTML or RTF body when plain body is empty
                if not body:
                    body = getattr(msg, "htmlBody", "") or ""
                if not body and hasattr(msg, "rtfBody") and msg.rtfBody:
                    try:
                        # optional: pip install striprtf
                        from striprtf.striprtf import rtf_to_text as rtf2txt
                        body = rtf2txt(
                            msg.rtfBody.decode(errors="ignore") if isinstance(msg.rtfBody, bytes) else msg.rtfBody
                        )
                    except Exception:
                        pass
        
                text = f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nDate: {date}\n\n{body}"
                return text[:MAX_OUT_CHARS]

            except Exception as e:
                logging.warning(f"MSG parse error: {e}")
                return None
            finally:
                try:
                    os.unlink(tmpname)
                except Exception:
                    pass
        # elif ext == '.msg' and extract_msg is not None:
        #     tmpname = "/tmp/tmps3_%s.msg" % (os.urandom(4).hex())
        #     with open(tmpname, "wb") as tmpf:
        #         tmpf.write(response['Body'].read())
        #     try:
        #         msg = extract_msg.Message(tmpname)
        #         text = f"Subject: {msg.subject}\nFrom: {msg.sender}\nTo: {msg.to}\nDate: {msg.date}\n\n{msg.body}"
        #         return text[:MAX_OUT_CHARS]
        #     except Exception as e:
        #         logging.warning(f"MSG parse error: {e}")
        #         return None
        #     finally:
        #         try: os.unlink(tmpname)
        #         except Exception: pass

        else:
            return ""
    except Exception as e:
        logging.warning(f"Failed to extract content from {bucket}:{key}: {e}")
        return None


def build_prompt(compliance_matrix):
    s = (
        "You are an expert compliance auditor and cyber expert. Based on your expert knowledge and using the table below for standards and regulated data fields as a reference,\n"
        "analyze the provided document (or data) for regulated data types and regulated data fields. When reviewing, pay close attention to the following and document your reasoning:\n"
        "- If the file contains or appears to contain only training material, instructional content with fake/sample/demo data, or random values (not real production data), clearly indicate this in your results and reasoning.\n"
        "- If regulated fields are present but there are no real (production or sensitive) values associated, rate the risk and data classification LOWER than if real production data is present.\n"
        "- If the data is explicitly marked as sample data, test records, or is clearly intended for demonstration or training, reduce the file's risk and sensitivity rating accordingly and document this in your explanation.\n"
        "- Only assign High or Medium risk and 'sensitive' classifications where there is strong evidence of actual regulated data populated in those fields.\n"
        "\n"
        "For every compliance standard with a match, create a results array entry containing:\n"
        "  - \"standard\": name of the standard\n"
        "  - \"jurisdiction\": the law's jurisdiction/region\n"
        "  - \"detected_fields\": list of regulated field names found in this document\n"
        "  - \"risk\": High, Medium, Low, or None, based on the actual data detected and the nature of that data (with sample/demo/fake data as Low or None)\n"
        "After reviewing all, output the following as ordered fields at the END of the JSON object:\n"
        "  - \"auditor_agent_view\": In 2-3 sentences, explain your reasoning (plain English, for auditors), specifically highlight if the content appears to be training/sample/demo, or fake/random data.\n"
        "  - \"likely_data_subject_area\": State data subject area, e.g., customer data, financial records, health information, etc.\n"
        "  - \"data_classification\": Classify Data, e.g., Highly Sensitive, Confidential, Internal Use, Public etc. Use a lower classification if the file only contains training/sample/demo data or regulated fields without actual sensitive values.\n"
        "  - \"overall_risk_rating\": string, High, Medium, Low, or None (the highest risk found overall; use Low or None if only sample/demo/empty data)\n"
        "  - \"hacker_interest\": State How Hacker or bad actors will potentially exploit this information if they get access to it. If the file only has demo/sample data, state 'Little to no value to hackers.'\n"
        "  - \"cyber_proposed_controls\": List proposed cyber controls leveraging ISO standard and NIST, based on data classification and risk rating as a cyber expert\n"
        "  - \"auditor_proposed_action\": string, short action phrase for remediation based on auditor expert knowledge, e.g., \"notify owner\", \"quarantine file\", \"no action required\" etc. For training/sample/demo data, recommend \"no action required\" or similar.\n"
        "Respond with a SINGLE pure JSON object only, in this format:\n"
        "{\n"
        "  \"results\": [\n"
        "     {\"standard\":\"\",\"jurisdiction\":\"\",\"detected_fields\":[],\"risk\":\"\",\"auditor_agent_view\":\"\"}, ...\n"
        "  ],\n"
        "  \"auditor_agent_view\": \"\",\n"
        "  \"likely_data_subject_area\": \"\",\n"
        "  \"data_classification\": \"\",\n"
        "  \"overall_risk_rating\": \"\",\n"
        "  \"hacker_interest\": \"\",\n"
        "  \"cyber_proposed_controls\":: \"\",\n"
        "  \"auditor_proposed_action\": \"\"\n"
        "}\n"
        "Here is your standards/fields matrix:\n"
    )
    for entry in compliance_matrix:
        s += (
            f"Standard: {entry['standard']} | Jurisdiction: {entry['jurisdiction']} | "
            f"Fields: {', '.join(entry['fields'])}\n"
        )
    return s

def gpt_classify_file(text, compliance_matrix):
    if not openai:
        logging.error("OpenAI module not installed or missing API key!")
        return "OpenAI not available"
    system_prompt = build_prompt(compliance_matrix)
    user = f"Classify the following file content for regulated data fields (see system context):\n\n{text[:MAX_OUT_CHARS]}"
    try:
        response = openai.chat.completions.create(
            model="gpt-4.1",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user}
            ],
            temperature=0.0,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        logging.error(f"GPT API error: {e}")
        return f"GPT API error: {e}"

def process_one_file(bucket, record, processed, compliance_matrix):
    key = record['key']
    if key in processed:
        return None  # Already done
    text = fetch_file_content(bucket, record)
    if text and text.strip():
        try:
            res = gpt_classify_file(text[:MAX_OUT_CHARS], compliance_matrix)
            outrec = dict(record)
            outrec['llm_response'] = res
            logging.info(f"Classified: [{bucket}] {record['file_name']}")
            return outrec
        except Exception as e:
            logging.error(f"Failed GPT for {key}: {e}")
            outrec = dict(record)
            outrec['llm_response'] = f"Error: {e}"
            return outrec
    else:
        logging.info(f"No usable text for file: [{bucket}] {record['file_name']}")
        return None

# ==== Main: For each bucket, call all matching per-bucket files ====
bucket_pattern = re.compile(r"([A-Za-z0-9._-]+)_s3_files_list_" + re.escape(config_id) + r"\.json$")
all_jsons = [fn for fn in os.listdir(config_dir) if bucket_pattern.match(fn)]

if not all_jsons:
    logging.warning(f"No bucket S3 file list JSONs found for config {config_id} in {config_dir}")
else:
    logging.info(f"Found {len(all_jsons)} S3 file lists: {all_jsons}")

for filelist_json in all_jsons:
    abs_filelist = os.path.join(config_dir, filelist_json)
    bucket = filelist_json.split('_s3_files_list_')[0]
    output_json = os.path.join(graph_dir, f"{bucket}_s3_output_{config_id}.json")

    records = load_json(abs_filelist, f"S3 file list for bucket {bucket}")

    if os.path.isfile(output_json):
        with open(output_json, "r", encoding="utf-8") as outf:
            output_list = json.load(outf)
        processed = set(r['key'] for r in output_list)
        logging.info(f"[{bucket}] Resuming: {len(processed)} files previously classified.")
    else:
        output_list = []
        processed = set()

    todo_records = [r for r in records if r['key'] not in processed]
    logging.info(f"[{bucket}] {len(todo_records)} files left to classify in parallel.")

    N = 12
    count = 0
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        futures = {executor.submit(process_one_file, bucket, rec, processed, compliance_matrix): rec for rec in todo_records}
        for future in as_completed(futures):
            outrec = future.result()
            if outrec:
                output_list.append(outrec)
                processed.add(outrec['key'])
                count += 1
                if count % N == 0:
                    with open(output_json, "w", encoding="utf-8") as outf:
                        json.dump(output_list, outf, indent=2)
                    logging.info(f"[{bucket}] {count} new this run, total {len(output_list)} files classified.")
    with open(output_json, "w", encoding="utf-8") as outf:
        json.dump(output_list, outf, indent=2)
    logging.info(f"[{bucket}] Completed. {len(output_list)} files classified.")