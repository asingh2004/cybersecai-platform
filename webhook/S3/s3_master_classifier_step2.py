import os
import sys
import json
import logging
import io
import re
import time
import base64
import tempfile
import subprocess
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed

from docx import Document
from PyPDF2 import PdfReader
import boto3

# Optional/extra filetype libraries
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None
try:
    from odf.opendocument import load as odf_load
    from odf.text import P
except ImportError:
    odf_load = None
    P = None
try:
    import extract_msg
except ImportError:
    extract_msg = None
try:
    from msg_parser import MsOxMessage as MsgParserMessage
except ImportError:
    MsgParserMessage = None
try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    import xlrd  # only supports .xls when xlrd<2.0
except ImportError:
    xlrd = None

import email
import email.policy

# OCR-related libs
try:
    from PIL import Image, ImageFile
    ImageFile.LOAD_TRUNCATED_IMAGES = True
except ImportError:
    Image = None
try:
    import pytesseract
except ImportError:
    pytesseract = None

# Optional PDF generation fallback (if LibreOffice not installed)
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import simpleSplit
except ImportError:
    canvas = None

# OpenAI (support both legacy and 1.x SDK)
try:
    from openai import OpenAI as OpenAIClientV1
except ImportError:
    OpenAIClientV1 = None
try:
    import openai as openai_legacy
except ImportError:
    openai_legacy = None

try:
    from botocore.exceptions import ClientError, BotoCoreError
except ImportError:
    ClientError = Exception
    BotoCoreError = Exception

logging.basicConfig(format="%(asctime)s %(levelname)s %(message)s", level=logging.INFO)

ALLOWED_EXT = (
    '.docx', '.txt', '.pdf', '.csv', '.json', '.md',
    '.xls', '.xlsx', '.ppt', '.pptx', '.rtf', '.odt', '.ods', '.eml', '.msg',
    # images for OCR
    '.jpg', '.jpeg', '.png', '.tif', '.tiff', '.gif', '.bmp', '.webp'
)

MAX_OUT_CHARS = 15000
MAX_WORKERS = 8
OCR_MAX_BYTES = 4_900_000  # safe limit under Textract Bytes 5MB
MAX_IMAGE_LONG_EDGE = int(os.getenv("OAI_MAX_IMAGE_LONG_EDGE", "2000"))  # px for OpenAI Vision downsizing
JPEG_QUALITY = int(os.getenv("OAI_JPEG_QUALITY", "85"))

# OpenAI model names
OAI_VISION_MODEL = os.getenv("OAI_VISION_MODEL", "gpt-4o-mini")
OAI_CLASSIFY_MODEL = os.getenv("OAI_CLASSIFY_MODEL", "gpt-4.1")

def abort(msg, code=1):
    logging.error(msg)
    sys.exit(code)

if len(sys.argv) < 2:
    abort("Usage: python3 s3_content_extract_classify.py <CONFIG_ID> [<search_root>]")

config_id = sys.argv[1]
search_root = sys.argv[2] if len(sys.argv) > 2 else '/home/cybersecai/htdocs/www.cybersecai.io/webhook/S3'

def find_config_folder(config_id, root):
    for rootdir, dirs, files in os.walk(root):
        if os.path.basename(rootdir) == config_id:
            return os.path.abspath(rootdir)
    abort(f"Could not find config folder '{config_id}' in '{root}'")

def load_json(path, desc):
    if not os.path.isfile(path):
        abort(f"{desc} not found: {path}")
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

config_dir = find_config_folder(config_id, search_root)
graph_dir  = os.path.join(config_dir, "graph")
os.makedirs(graph_dir, exist_ok=True)

secrets_json = os.path.join(config_dir, f"{config_id}_secrets.json")
compliance_matrix = load_json(os.path.join(config_dir, "compliance_matrix.json"), "Compliance matrix")
s3config = load_json(secrets_json, "S3 Secrets file")
AWS_ACCESS_KEY = s3config['AWS_ACCESS_KEY']
AWS_SECRET_KEY = s3config['AWS_SECRET_KEY']
region = s3config.get('AWS_REGION', 'us-east-1')

s3 = boto3.client(
    "s3",
    aws_access_key_id=AWS_ACCESS_KEY,
    aws_secret_access_key=AWS_SECRET_KEY,
    region_name=region
)

textract = boto3.client(
    "textract",
    aws_access_key_id=AWS_ACCESS_KEY,
    aws_secret_access_key=AWS_SECRET_KEY,
    region_name=region
)

def _retry(func, *args, retries=3, backoff=1.5, **kwargs):
    last_exc = None
    for attempt in range(retries):
        try:
            return func(*args, **kwargs)
        except (ClientError, BotoCoreError, Exception) as e:
            last_exc = e
            if attempt < retries - 1:
                sleep_for = backoff ** attempt
                logging.warning(f"{func.__name__} failed (attempt {attempt+1}/{retries}): {e}. Retrying in {sleep_for:.1f}s")
                time.sleep(sleep_for)
            else:
                logging.error(f"{func.__name__} failed after {retries} attempts: {e}")
    raise last_exc

def which(bin_name):
    return shutil.which(bin_name)

HAS_SOFFICE = which("soffice") is not None
HAS_MSGCONVERT = which("msgconvert") is not None  # from libemail-outlook-message-perl (Linux)

# ---------------------- OpenAI helpers ---------------------- #
def get_openai_client():
    api_key = os.getenv("OPENAI_API_KEY") or os.getenv("OPENAI_APIKEY") or os.getenv("OPENAI_KEY")
    if not api_key:
        return None, None
    if OpenAIClientV1 is not None:
        return OpenAIClientV1(api_key=api_key), "v1"
    if openai_legacy is not None:
        openai_legacy.api_key = api_key
        return openai_legacy, "legacy"
    return None, None

def prepare_image_for_vision(image_bytes):
    if Image is None:
        return image_bytes, "image/png"
    try:
        with Image.open(io.BytesIO(image_bytes)) as im:
            if im.mode not in ("RGB", "L"):
                im = im.convert("RGB")
            w, h = im.size
            long_edge = max(w, h)
            if long_edge > MAX_IMAGE_LONG_EDGE:
                scale = MAX_IMAGE_LONG_EDGE / float(long_edge)
                im = im.resize((max(1, int(w*scale)), max(1, int(h*scale))))
            buf = io.BytesIO()
            fmt = "JPEG" if im.mode in ("RGB", "L") else "PNG"
            if fmt == "JPEG":
                im.save(buf, format=fmt, quality=JPEG_QUALITY, optimize=True)
                mime = "image/jpeg"
            else:
                im.save(buf, format="PNG", optimize=True)
                mime = "image/png"
            return buf.getvalue(), mime
    except Exception as e:
        logging.warning(f"prepare_image_for_vision failed, using original: {e}")
        return image_bytes, "application/octet-stream"

def openai_vision_ocr(image_bytes):
    client, mode = get_openai_client()
    if client is None:
        return None
    try:
        prepped, mime = prepare_image_for_vision(image_bytes)
        data_url = f"data:{mime};base64,{base64.b64encode(prepped).decode('ascii')}"
        user_content = [
            {"type": "text", "text": "Extract all visible text from the image. Return only the raw text."},
            {"type": "image_url", "image_url": {"url": data_url}},
        ]
        if mode == "v1":
            resp = client.chat.completions.create(
                model=OAI_VISION_MODEL,
                messages=[{"role": "user", "content": user_content}],
                temperature=0,
            )
            return (resp.choices[0].message.content or "").strip()
        else:
            resp = client.ChatCompletion.create(
                model=OAI_VISION_MODEL,
                messages=[{"role": "user", "content": user_content}],
                temperature=0,
            )
            return (resp["choices"][0]["message"]["content"] or "").strip()
    except Exception as e:
        logging.warning(f"OpenAI Vision OCR error: {e}")
        return None

# ---------------------- Textract helpers ---------------------- #
def textract_detect_lines_from_s3(bucket, key):
    resp = _retry(textract.detect_document_text, Document={'S3Object': {'Bucket': bucket, 'Name': key}})
    lines = [b['Text'] for b in resp.get('Blocks', []) if b.get('BlockType') == 'LINE' and b.get('Text')]
    return "\n".join(lines)

def textract_detect_lines_from_bytes(img_bytes):
    resp = _retry(textract.detect_document_text, Document={'Bytes': img_bytes})
    lines = [b['Text'] for b in resp.get('Blocks', []) if b.get('BlockType') == 'LINE' and b.get('Text')]
    return "\n".join(lines)

def pil_convert_to_png_bytes(image_bytes):
    if Image is None:
        return None
    try:
        with Image.open(io.BytesIO(image_bytes)) as im:
            with io.BytesIO() as out:
                im = im.convert("RGB")
                im.save(out, format="PNG", optimize=True)
                return out.getvalue()
    except Exception as e:
        logging.warning(f"PIL conversion to PNG failed: {e}")
        return None

def pytesseract_ocr(image_bytes):
    if Image is None or pytesseract is None:
        return None
    try:
        with Image.open(io.BytesIO(image_bytes)) as im:
            return pytesseract.image_to_string(im)
    except Exception as e:
        logging.warning(f"pytesseract OCR failed: {e}")
        return None

def ocr_image(bucket, record, ext):
    key = record['key']
    direct_textract_ext = ('.jpg', '.jpeg', '.png', '.tif', '.tiff', '.pdf')
    try:
        if ext in direct_textract_ext:
            try:
                text = textract_detect_lines_from_s3(bucket, key)
                if text and text.strip():
                    return text[:MAX_OUT_CHARS]
            except Exception as e:
                logging.warning(f"Textract S3 OCR failed: {e}")

        response = _retry(s3.get_object, Bucket=bucket, Key=key)
        blob = response['Body'].read(OCR_MAX_BYTES + 200_000)

        try:
            png_bytes = pil_convert_to_png_bytes(blob) if ext not in ('.jpg', '.jpeg', '.png', '.tif', '.tiff') else None
            if png_bytes and len(png_bytes) <= OCR_MAX_BYTES:
                text = textract_detect_lines_from_bytes(png_bytes)
                if text and text.strip():
                    return text[:MAX_OUT_CHARS]
        except Exception as e:
            logging.warning(f"Textract Bytes OCR failed: {e}")

        try:
            text = openai_vision_ocr(blob)
            if text and text.strip():
                return text[:MAX_OUT_CHARS]
        except Exception as e:
            logging.warning(f"OpenAI Vision OCR fallback failed: {e}")

        text = pytesseract_ocr(blob)
        return text[:MAX_OUT_CHARS] if text else None

    except Exception as e:
        logging.warning(f"OCR failed for {bucket}:{key} ({ext}): {e}")
        try:
            response = _retry(s3.get_object, Bucket=bucket, Key=key)
            blob = response['Body'].read(OCR_MAX_BYTES)
            text = pytesseract_ocr(blob)
            return text[:MAX_OUT_CHARS] if text else None
        except Exception as e2:
            logging.warning(f"OCR fallback failed for {bucket}:{key}: {e2}")
            return None

# ---------------------- Conversion helpers ---------------------- #
def libreoffice_convert(input_path, outdir, target):
    """
    Generic LO headless convert. target examples:
      - 'pdf:writer_pdf_Export' (Writer)
      - 'pdf:calc_pdf_Export' (Calc)
      - 'pdf:impress_pdf_Export' (Impress)
      - 'pdf' (auto)
      - 'xlsx' / 'docx' / 'csv' etc.
    """
    if not HAS_SOFFICE:
        return None
    cmd = ["soffice", "--headless", "--norestore", "--convert-to", target, "--outdir", outdir, input_path]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=240)
        base = os.path.splitext(os.path.basename(input_path))[0]
        # Infer output extension from target
        ext = target.split(":")[0]
        if ext in ("pdf", "docx", "xlsx", "csv", "odt", "ods"):
            outpath = os.path.join(outdir, base + "." + ext)
            if os.path.isfile(outpath):
                return outpath
        # Fallback: find by scanning outdir for files with same base
        for fn in os.listdir(outdir):
            if fn.startswith(base + "."):
                candidate = os.path.join(outdir, fn)
                if os.path.isfile(candidate):
                    return candidate
    except Exception as e:
        logging.warning(f"LibreOffice convert failed ({target}): {e}")
    return None

def libreoffice_convert_to_pdf(input_path, work_dir=None, kind_hint=None):
    """
    Convert document to PDF using headless LibreOffice.
    kind_hint: 'writer'|'calc'|'impress' to choose better PDF filter.
    """
    outdir = work_dir or os.path.dirname(input_path)
    target = "pdf"
    if kind_hint == "writer":
        target = "pdf:writer_pdf_Export"
    elif kind_hint == "calc":
        target = "pdf:calc_pdf_Export"
    elif kind_hint == "impress":
        target = "pdf:impress_pdf_Export"
    return libreoffice_convert(input_path, outdir, target)

def extract_text_from_pdf_bytes(pdf_bytes, ocr_if_empty=True):
    text = ""
    try:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            text += (page.extract_text() or "") + "\n"
    except Exception as e:
        logging.warning(f"PDF parse error: {e}")
    if ocr_if_empty and len(text.strip()) < 100:
        # Optional: Could add pdf2image->OCR here if desired
        pass
    return text

# ---------------------- .msg/.eml helpers ---------------------- #
def parse_msg_with_extract_msg(path):
    if extract_msg is None:
        return None
    try:
        msg = extract_msg.Message(path)
        subject = msg.subject or ""
        sender = getattr(msg, "sender", "") or getattr(msg, "sender_email", "")
        to = msg.to or ""
        date = msg.date or ""
        body = msg.body or ""
        if not body:
            body = getattr(msg, "htmlBody", "") or ""
        if not body and hasattr(msg, "rtfBody") and msg.rtfBody:
            try:
                if isinstance(msg.rtfBody, bytes):
                    body = rtf_to_text(msg.rtfBody.decode(errors="ignore")) if rtf_to_text else ""
                else:
                    body = rtf_to_text(msg.rtfBody) if rtf_to_text else ""
            except Exception:
                pass
        text = f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nDate: {date}\n\n{body}"
        try:
            msg.close()
        except Exception:
            pass
        return text
    except Exception as e:
        logging.warning(f"extract_msg parse error: {e}")
        return None

def parse_msg_with_msg_parser(path):
    if MsgParserMessage is None:
        return None
    try:
        m = MsgParserMessage(path)
        m_message = m.get_email_data()
        subject = m_message.get("subject", "") or ""
        sender = m_message.get("sender", "") or ""
        to = ", ".join(m_message.get("to", []) or [])
        date = m_message.get("date", "") or ""
        body = (m_message.get("body", "") or "") + "\n" + (m_message.get("body_html", "") or "")
        text = f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nDate: {date}\n\n{body}"
        return text
    except Exception as e:
        logging.warning(f"msg_parser parse error: {e}")
        return None

def convert_msg_to_eml_with_msgconvert(msg_path, out_dir):
    if not HAS_MSGCONVERT:
        return None
    try:
        env = os.environ.copy()
        cwd = out_dir
        res = subprocess.run(["msgconvert", msg_path], cwd=cwd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)
        if res.returncode != 0:
            logging.warning(f"msgconvert failed: {res.stderr.decode(errors='ignore')}")
            return None
        # Find newly created .eml with same base
        base = os.path.splitext(os.path.basename(msg_path))[0]
        for fn in os.listdir(out_dir):
            if fn.lower().startswith(base.lower()) and fn.lower().endswith(".eml"):
                return os.path.join(out_dir, fn)
        # If not found, look for any .eml created latest
        emls = [os.path.join(out_dir, fn) for fn in os.listdir(out_dir) if fn.lower().endswith(".eml")]
        if emls:
            return sorted(emls, key=lambda p: os.path.getmtime(p), reverse=True)[0]
    except Exception as e:
        logging.warning(f"msgconvert exception: {e}")
    return None

def parse_eml_bytes(content):
    try:
        msg = email.message_from_bytes(content, policy=email.policy.default)
        text_parts = []
        if msg['Subject']:
            text_parts.append("Subject: " + str(msg['Subject']))
        if msg['From']:
            text_parts.append("From: " + str(msg['From']))
        if msg['To']:
            text_parts.append("To: " + str(msg['To']))
        if msg['Date']:
            text_parts.append("Date: " + str(msg['Date']))
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                disp = (part.get("Content-Disposition") or "").lower()
                if ctype == "text/plain" and "attachment" not in disp:
                    try:
                        body += part.get_content()
                    except Exception:
                        try:
                            body += (part.get_payload(decode=True) or b"").decode(errors="ignore")
                        except Exception:
                            pass
        else:
            try:
                body = msg.get_content()
            except Exception:
                body = (msg.get_payload(decode=True) or b"").decode(errors="ignore")
        text_parts.append(str(body))
        return "\n".join(text_parts)
    except Exception as e:
        logging.warning(f"EML parse error: {e}")
        return None

def simple_extract_msg_text(path):
    """
    Extracts text from .msg using Python libraries and msgconvert->eml parse,
    without trying to convert to PDF. Used as a fallback for PDF generation.
    """
    # 1) Native libraries
    text = parse_msg_with_extract_msg(path)
    if text and text.strip():
        return text[:MAX_OUT_CHARS]
    text = parse_msg_with_msg_parser(path)
    if text and text.strip():
        return text[:MAX_OUT_CHARS]
    # 2) Convert to EML and parse
    with tempfile.TemporaryDirectory() as td:
        eml_path = convert_msg_to_eml_with_msgconvert(path, td)
        if eml_path and os.path.isfile(eml_path):
            try:
                with open(eml_path, "rb") as f:
                    content = f.read()
                text = parse_eml_bytes(content)
                if text and text.strip():
                    return text[:MAX_OUT_CHARS]
            except Exception as e:
                logging.warning(f"EML read after msgconvert failed: {e}")
    return None

def html_escape(s):
    try:
        import html
        return html.escape(s)
    except Exception:
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def write_text_as_html_file(text, out_dir, base="email_export"):
    html_content = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>Email Export</title>
<style>
body {{ font-family: Arial, sans-serif; white-space: pre-wrap; }}
.header {{ font-weight: bold; margin-bottom: 10px; }}
</style>
</head>
<body>
{html_escape(text)}
</body>
</html>
"""
    html_path = os.path.join(out_dir, f"{base}.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    return html_path

def reportlab_write_pdf(text, out_pdf):
    if canvas is None:
        return False
    try:
        c = canvas.Canvas(out_pdf, pagesize=letter)
        width, height = letter
        margin = 36  # half inch
        max_width = width - 2*margin
        y = height - margin
        # Use a monospaced font to preserve simple formatting
        c.setFont("Helvetica", 10)
        lines = text.splitlines() or [text]
        for line in lines:
            # Wrap long lines
            wrapped = simpleSplit(line, "Helvetica", 10, max_width)
            for wline in wrapped:
                if y <= margin:
                    c.showPage()
                    c.setFont("Helvetica", 10)
                    y = height - margin
                c.drawString(margin, y, wline)
                y -= 12
        c.save()
        return True
    except Exception as e:
        logging.warning(f"ReportLab PDF generation failed: {e}")
        return False

def convert_msg_to_pdf(msg_path, work_dir=None):
    """
    Ensures a .pdf is produced from a .msg file using robust fallbacks.
    Returns path to PDF or None.
    """
    td_created = False
    if work_dir is None:
        work_dir = tempfile.mkdtemp()
        td_created = True
    pdf_path = None
    try:
        # 1) msgconvert -> .eml -> LibreOffice PDF
        eml_path = convert_msg_to_eml_with_msgconvert(msg_path, work_dir)
        if eml_path and os.path.isfile(eml_path):
            pdf_path = libreoffice_convert_to_pdf(eml_path, work_dir=work_dir, kind_hint="writer")
            if pdf_path and os.path.isfile(pdf_path):
                return pdf_path

        # 2) Try LibreOffice convert .msg directly
        pdf_path = libreoffice_convert_to_pdf(msg_path, work_dir=work_dir, kind_hint="writer")
        if pdf_path and os.path.isfile(pdf_path):
            return pdf_path

        # 3) Extract text and then convert HTML->PDF via LibreOffice or ReportLab
        text = simple_extract_msg_text(msg_path) or ""
        if text.strip():
            html_path = write_text_as_html_file(text, work_dir, base="email_from_msg")
            # 3a) LibreOffice from HTML to PDF
            pdf_path = libreoffice_convert_to_pdf(html_path, work_dir=work_dir, kind_hint="writer")
            if pdf_path and os.path.isfile(pdf_path):
                return pdf_path
            # 3b) ReportLab direct PDF as ultimate fallback
            rl_pdf = os.path.join(work_dir, "email_from_msg_reportlab.pdf")
            if reportlab_write_pdf(text, rl_pdf):
                return rl_pdf

        logging.warning("convert_msg_to_pdf: All PDF conversion paths failed.")
        return None
    finally:
        if td_created:
            # Don't delete temp dir immediately if we return path inside it, so only clean if None
            if pdf_path is None:
                try:
                    for fn in os.listdir(work_dir):
                        os.unlink(os.path.join(work_dir, fn))
                    os.rmdir(work_dir)
                except Exception:
                    pass

# Legacy function kept for reference, but not used for classification anymore
def parse_msg_from_file(path):
    """
    Deprecated in this flow. Kept for reference. We now convert .msg to PDF then extract text.
    """
    text = parse_msg_with_extract_msg(path)
    if text and text.strip():
        return text[:MAX_OUT_CHARS]
    text = parse_msg_with_msg_parser(path)
    if text and text.strip():
        return text[:MAX_OUT_CHARS]

    with tempfile.TemporaryDirectory() as td:
        eml_path = convert_msg_to_eml_with_msgconvert(path, td)
        if eml_path and os.path.isfile(eml_path):
            try:
                with open(eml_path, "rb") as f:
                    content = f.read()
                text = parse_eml_bytes(content)
                if text and text.strip():
                    return text[:MAX_OUT_CHARS]
            except Exception as e:
                logging.warning(f"EML read after msgconvert failed: {e}")
        if eml_path and HAS_SOFFICE:
            pdf_path = libreoffice_convert_to_pdf(eml_path, work_dir=td, kind_hint="writer")
            if pdf_path and os.path.isfile(pdf_path):
                with open(pdf_path, "rb") as pf:
                    pdf_bytes = pf.read()
                text = extract_text_from_pdf_bytes(pdf_bytes, ocr_if_empty=True)
                if text and text.strip():
                    return text[:MAX_OUT_CHARS]

    if HAS_SOFFICE:
        td2 = tempfile.mkdtemp()
        try:
            pdf_path = libreoffice_convert_to_pdf(path, work_dir=td2, kind_hint="writer")
            if pdf_path and os.path.isfile(pdf_path):
                with open(pdf_path, "rb") as pf:
                    pdf_bytes = pf.read()
                text = extract_text_from_pdf_bytes(pdf_bytes, ocr_if_empty=True)
                if text and text.strip():
                    return text[:MAX_OUT_CHARS]
        finally:
            try:
                for fn in os.listdir(td2):
                    os.unlink(os.path.join(td2, fn))
                os.rmdir(td2)
            except Exception:
                pass

    return None

# ---------------------- Excel helpers ---------------------- #
def excel_to_text_with_openpyxl(xlsx_path, max_rows=1000):
    if openpyxl is None:
        return None
    try:
        wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
        pieces = []
        for ws in wb.worksheets:
            rows = []
            rcount = 0
            for row in ws.iter_rows(values_only=True):
                rcount += 1
                if rcount > max_rows:
                    break
                vals = [(str(v) if v is not None else "") for v in row]
                rows.append(",".join(vals))
            if rows:
                pieces.append(f"=== Sheet: {ws.title} ===\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(pieces)
    except Exception as e:
        logging.warning(f"openpyxl extract error: {e}")
        return None

def excel_xls_to_text_with_xlrd(xls_path, max_rows=1000):
    if xlrd is None:
        return None
    try:
        book = xlrd.open_workbook(xls_path, on_demand=True)
        pieces = []
        for sheet_name in book.sheet_names():
            sh = book.sheet_by_name(sheet_name)
            rows = []
            for r in range(min(sh.nrows, max_rows)):
                vals = []
                for c in range(sh.ncols):
                    cell = sh.cell(r, c)
                    if cell.ctype == xlrd.XL_CELL_DATE:
                        try:
                            from datetime import datetime, timedelta
                            dt = xlrd.xldate_as_tuple(cell.value, book.datemode)
                            vals.append(str(datetime(*dt)))
                        except Exception:
                            vals.append(str(cell.value))
                    else:
                        vals.append(str(cell.value) if cell.value is not None else "")
                rows.append(",".join(vals))
            if rows:
                pieces.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
        book.release_resources()
        return "\n\n".join(pieces)
    except Exception as e:
        logging.warning(f"xlrd extract error: {e}")
        return None

def excel_bytes_to_text(blob, ext):
    # Try local parsers first
    with tempfile.TemporaryDirectory() as td:
        in_path = os.path.join(td, f"input{ext}")
        with open(in_path, "wb") as f:
            f.write(blob)

        if ext == ".xlsx":
            txt = excel_to_text_with_openpyxl(in_path)
            if txt and txt.strip():
                return txt

            # Fallback: pandas if available
            if pd is not None:
                try:
                    xls = pd.ExcelFile(in_path, engine='openpyxl')
                    pieces = []
                    for sheet in xls.sheet_names:
                        try:
                            df = pd.read_excel(xls, sheet_name=sheet, dtype=str, nrows=1000, na_filter=False)
                            if not df.empty:
                                csv_text = df.to_csv(index=False)
                                if csv_text.strip():
                                    pieces.append(f"=== Sheet: {sheet} ===\n{csv_text}")
                        except Exception as se:
                            logging.warning(f"Pandas Excel sheet parse error ({sheet}): {se}")
                    if pieces:
                        return "\n\n".join(pieces)
                except Exception as e:
                    logging.warning(f"Pandas Excel xlsx parse error: {e}")

        else:  # .xls
            txt = excel_xls_to_text_with_xlrd(in_path)
            if txt and txt.strip():
                return txt

            # Fallback: pandas if available and xlrd supports .xls
            if pd is not None:
                try:
                    xls = pd.ExcelFile(in_path)  # let pandas choose engine
                    pieces = []
                    for sheet in xls.sheet_names:
                        try:
                            df = pd.read_excel(xls, sheet_name=sheet, dtype=str, nrows=1000, na_filter=False)
                            if not df.empty:
                                csv_text = df.to_csv(index=False)
                                if csv_text.strip():
                                    pieces.append(f"=== Sheet: {sheet} ===\n{csv_text}")
                        except Exception as se:
                            logging.warning(f"Pandas Excel sheet parse error ({sheet}): {se}")
                    if pieces:
                        return "\n\n".join(pieces)
                except Exception as e:
                    logging.warning(f"Pandas Excel xls parse error: {e}")

        # Fallback via LibreOffice:
        # 1) Try convert to .xlsx then use openpyxl
        if HAS_SOFFICE:
            xlsx_path = libreoffice_convert(in_path, td, "xlsx")
            if xlsx_path and os.path.isfile(xlsx_path):
                txt = excel_to_text_with_openpyxl(xlsx_path)
                if txt and txt.strip():
                    return txt
            # 2) Try convert directly to PDF then extract text (last resort)
            pdf_path = libreoffice_convert_to_pdf(in_path, work_dir=td, kind_hint="calc")
            if pdf_path and os.path.isfile(pdf_path):
                with open(pdf_path, "rb") as pf:
                    pdf_bytes = pf.read()
                text = extract_text_from_pdf_bytes(pdf_bytes, ocr_if_empty=True)
                if text and text.strip():
                    return text

    return None

# ---------------------- File content extraction ---------------------- #
def parse_pptx_from_file(path):
    if Presentation is not None:
        try:
            prs = Presentation(path)
            lines = []
            for slide in prs.slides:
                if getattr(slide.shapes, "title", None) and getattr(slide.shapes.title, "text", None):
                    lines.append(slide.shapes.title.text)
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and getattr(shape, "text", None):
                        lines.append(shape.text)
                    if hasattr(shape, "table") and shape.table is not None:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    lines.append(cell.text)
            combined = "\n".join(lines).strip()
            if combined:
                return combined[:MAX_OUT_CHARS]
        except Exception as e:
            logging.warning(f"python-pptx parse error: {e}")
    pdf_path = libreoffice_convert_to_pdf(path, work_dir=os.path.dirname(path), kind_hint="impress")
    if pdf_path and os.path.isfile(pdf_path):
        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()
        text = extract_text_from_pdf_bytes(pdf_bytes)
        try:
            os.unlink(pdf_path)
        except Exception:
            pass
        return text[:MAX_OUT_CHARS] if text else None
    return None

def parse_ppt_from_bytes(blob, ext):
    with tempfile.TemporaryDirectory() as td:
        in_path = os.path.join(td, f"input{ext}")
        with open(in_path, "wb") as f:
            f.write(blob)
        pdf_path = libreoffice_convert_to_pdf(in_path, work_dir=td, kind_hint="impress")
        if pdf_path and os.path.isfile(pdf_path):
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()
            text = extract_text_from_pdf_bytes(pdf_bytes)
            return text[:MAX_OUT_CHARS] if text else None
    return None

def fetch_file_content(bucket, record):
    key = record['key']
    ext = os.path.splitext(record['file_name'])[1].lower()
    if ext not in ALLOWED_EXT or record.get('file_type','') == 'folder':
        return None
    try:
        if ext in ('.jpg', '.jpeg', '.png', '.tif', '.tiff', '.gif', '.bmp', '.webp'):
            return ocr_image(bucket, record, ext)

        response = s3.get_object(Bucket=bucket, Key=key)

        if ext in ('.txt', '.csv', '.json', '.md'):
            content = response['Body'].read(MAX_OUT_CHARS*2)
            try:
                return content.decode('utf-8')[:MAX_OUT_CHARS]
            except Exception:
                return content.decode('latin-1', errors="replace")[:MAX_OUT_CHARS]

        elif ext == '.docx':
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmpf:
                tmpf.write(response['Body'].read())
                tmpname = tmpf.name
            try:
                doc = Document(tmpname)
                return "\n".join(p.text for p in doc.paragraphs)[:MAX_OUT_CHARS]
            finally:
                try: os.unlink(tmpname)
                except Exception: pass

        elif ext == '.pdf':
            content = response['Body'].read(15_000_000)
            text = extract_text_from_pdf_bytes(content, ocr_if_empty=True)
            if len((text or "").strip()) < 100:
                try:
                    text2 = textract_detect_lines_from_s3(bucket, key)
                    if text2 and len(text2.strip()) > len((text or "").strip()):
                        text = text2
                except Exception as e:
                    logging.warning(f"Textract OCR for PDF failed: {e}")
            return text[:MAX_OUT_CHARS] if text else None

        elif ext in ('.xls', '.xlsx'):
            blob = response['Body'].read()
            text = excel_bytes_to_text(blob, ext)
            return text[:MAX_OUT_CHARS] if text else None

        elif ext in ('.pptx', '.ppt'):
            blob = response['Body'].read()
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmpf:
                tmpf.write(blob)
                inpath = tmpf.name
            try:
                if ext == '.pptx':
                    text = parse_pptx_from_file(inpath)
                else:
                    text = parse_ppt_from_bytes(blob, ext)
                if (not text) or len(text.strip()) < 10:
                    pdf_path = libreoffice_convert_to_pdf(inpath, kind_hint="impress")
                    if pdf_path and os.path.isfile(pdf_path):
                        with open(pdf_path, "rb") as pf:
                            pdf_bytes = pf.read()
                        text = extract_text_from_pdf_bytes(pdf_bytes)
                return text[:MAX_OUT_CHARS] if text else None
            finally:
                try: os.unlink(inpath)
                except Exception: pass

        elif ext == '.rtf' and rtf_to_text is not None:
            content = response['Body'].read()
            try:
                return rtf_to_text(content.decode(errors="ignore"))[:MAX_OUT_CHARS]
            except Exception as e:
                logging.warning(f"RTF parse error: {e}")
                return None

        elif ext in ('.odt', '.ods') and odf_load is not None and P is not None:
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmpf:
                tmpf.write(response['Body'].read())
                tmpname = tmpf.name
            try:
                doc = odf_load(tmpname)
                paragraphs = []
                for elem in doc.getElementsByType(P):
                    try:
                        paragraphs.append(elem.firstChild.data if elem.firstChild else "")
                    except Exception:
                        paragraphs.append(str(elem))
                combined = "\n".join(paragraphs)
                return combined[:MAX_OUT_CHARS] if combined else None
            except Exception as e:
                logging.warning(f"ODF parse error: {e}")
                return None
            finally:
                try: os.unlink(tmpname)
                except Exception: pass

        elif ext == '.eml':
            content = response['Body'].read(2_000_000)
            text = parse_eml_bytes(content)
            return text[:MAX_OUT_CHARS] if text else None

        elif ext == '.msg':
            # New behavior: convert .msg to .pdf first, then extract text from PDF
            blob = response['Body'].read()
            with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as tmpf:
                tmpf.write(blob)
                tmpname = tmpf.name
            try:
                pdf_path = convert_msg_to_pdf(tmpname)
                text = None
                if pdf_path and os.path.isfile(pdf_path):
                    try:
                        with open(pdf_path, "rb") as pf:
                            pdf_bytes = pf.read()
                        text = extract_text_from_pdf_bytes(pdf_bytes, ocr_if_empty=True)
                        if (not text) or len(text.strip()) < 50:
                            # As a secondary attempt, parse msg text directly and append
                            extra_text = simple_extract_msg_text(tmpname)
                            if extra_text:
                                text = (text or "") + "\n\n" + extra_text
                    except Exception as e:
                        logging.warning(f"Reading/extracting from generated PDF failed: {e}")
                else:
                    logging.warning("MSG->PDF conversion failed; falling back to direct text extraction.")
                    text = simple_extract_msg_text(tmpname)

                return text[:MAX_OUT_CHARS] if text else None
            finally:
                try: os.unlink(tmpname)
                except Exception: pass

        else:
            return ""
    except Exception as e:
        logging.warning(f"Failed to extract content from {bucket}:{key}: {e}")
        return None

# ---------------------- Classification prompt ---------------------- #
def build_prompt(compliance_matrix):
    s = (
        "You are an expert compliance auditor and cyber expert. Based on your expert knowledge and using the table below for standards and regulated data fields as a reference,\n"
        "analyze the provided document (or data) for regulated data types and regulated data fields. When reviewing, pay close attention to the following and document your reasoning:\n"
        "- If the file contains or appears to contain only training material, instructional content with fake/sample/demo data, or random values (not real production data), clearly indicate this in your results and reasoning.\n"
        "- If regulated fields are present but there are no real (production or sensitive) values associated, rate the risk and data classification LOWER than if real production data is present.\n"
        "- If the data is explicitly marked as sample data, test records, or is clearly intended for demonstration or training, reduce the file's risk and sensitivity rating accordingly and document this in your explanation.\n"
        "- Only assign High or Medium risk and 'sensitive' classifications where there is strong evidence of actual regulated data populated in those fields.\n"
        "\n"
        "For every compliance standard with a match, create a results array entry containing:\n"
        "  - \"standard\": name of the standard\n"
        "  - \"jurisdiction\": the law's jurisdiction/region\n"
        "  - \"detected_fields\": list of regulated field names found in this document\n"
        "  - \"risk\": High, Medium, Low, or None, based on the actual data detected and the nature of that data (with sample/demo/fake data as Low or None)\n"
        "After reviewing all, output the following as ordered fields at the END of the JSON object:\n"
        "  - \"auditor_agent_view\": In 2-3 sentences, explain your reasoning (plain English, for auditors), specifically highlight if the content appears to be training/sample/demo, or fake/random data.\n"
        "  - \"likely_data_subject_area\": State data subject area, e.g., customer data, financial records, health information, etc.\n"
        "  - \"data_classification\": Classify Data, e.g., Highly Sensitive, Confidential, Internal Use, Public etc. Use a lower classification if the file only contains training/sample/demo data or regulated fields without actual sensitive values.\n"
        "  - \"overall_risk_rating\": string, High, Medium, Low, or None (the highest risk found overall; use Low or None if only sample/demo/empty data)\n"
        "  - \"hacker_interest\": State How Hacker or bad actors will potentially exploit this information if they get access to it. If the file only has demo/sample data, state 'Little to no value to hackers.'\n"
        "  - \"cyber_proposed_controls\": List proposed cyber controls leveraging ISO standard and NIST, based on data classification and risk rating as a cyber expert\n"
        "  - \"auditor_proposed_action\": string, short action phrase for remediation based on auditor expert knowledge, e.g., \"notify owner\", \"quarantine file\", \"no action required\" etc. For training/sample/demo data, recommend \"no action required\" or similar.\n"
        "Respond with a SINGLE pure JSON object only, in this format:\n"
        "{\n"
        "  \"results\": [\n"
        "     {\"standard\":\"\",\"jurisdiction\":\"\",\"detected_fields\":[],\"risk\":\"\",\"auditor_agent_view\":\"\"}, ...\n"
        "  ],\n"
        "  \"auditor_agent_view\": \"\",\n"
        "  \"likely_data_subject_area\": \"\",\n"
        "  \"data_classification\": \"\",\n"
        "  \"overall_risk_rating\": \"\",\n"
        "  \"hacker_interest\": \"\",\n"
        "  \"cyber_proposed_controls\": \"\",\n"
        "  \"auditor_proposed_action\": \"\"\n"
        "}\n"
        "Here is your standards/fields matrix:\n"
    )
    for entry in compliance_matrix:
        s += (
            f"Standard: {entry['standard']} | Jurisdiction: {entry['jurisdiction']} | "
            f"Fields: {', '.join(entry['fields'])}\n"
        )
    return s

def gpt_classify_file(text, compliance_matrix):
    client, mode = get_openai_client()
    if client is None:
        logging.error("OpenAI SDK not installed or OPENAI_API_KEY not set; cannot classify.")
        return "OpenAI not available"
    system_prompt = build_prompt(compliance_matrix)
    user = f"Classify the following file content for regulated data fields (see system context):\n\n{text[:MAX_OUT_CHARS]}"
    try:
        if mode == "v1":
            response = client.chat.completions.create(
                model=OAI_CLASSIFY_MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user}
                ],
                temperature=0.0,
            )
            return (response.choices[0].message.content or "").strip()
        else:
            response = client.ChatCompletion.create(
                model=OAI_CLASSIFY_MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user}
                ],
                temperature=0.0,
            )
            return (response["choices"][0]["message"]["content"] or "").strip()
    except Exception as e:
        logging.error(f"GPT API error: {e}")
        return f"GPT API error: {e}"

# ---------------------- Orchestration ---------------------- #
def process_one_file(bucket, record, processed, compliance_matrix):
    key = record['key']
    if key in processed:
        return None
    text = fetch_file_content(bucket, record)
    if text and text.strip():
        try:
            res = gpt_classify_file(text[:MAX_OUT_CHARS], compliance_matrix)
            outrec = dict(record)
            outrec['llm_response'] = res
            logging.info(f"Classified: [{bucket}] {record['file_name']}")
            return outrec
        except Exception as e:
            logging.error(f"Failed GPT for {key}: {e}")
            outrec = dict(record)
            outrec['llm_response'] = f"Error: {e}"
            return outrec
    else:
        logging.info(f"No usable text for file: [{bucket}] {record['file_name']}")
        return None

def which_bucket_jsons(config_dir, config_id):
    bucket_pattern = re.compile(r"([A-Za-z0-9._-]+)_s3_files_list_" + re.escape(config_id) + r"\.json$")
    return [fn for fn in os.listdir(config_dir) if bucket_pattern.match(fn)], bucket_pattern

all_jsons, bucket_pattern = which_bucket_jsons(config_dir, config_id)

if not all_jsons:
    logging.warning(f"No S3 file list JSONs found for config {config_id} in {config_dir}")
else:
    logging.info(f"Found {len(all_jsons)} S3 file lists: {all_jsons}")

for filelist_json in all_jsons:
    abs_filelist = os.path.join(config_dir, filelist_json)
    bucket = filelist_json.split('_s3_files_list_')[0]
    output_json = os.path.join(graph_dir, f"{bucket}_s3_output_{config_id}.json")
    records = load_json(abs_filelist, f"S3 file list for bucket {bucket}")

    if os.path.isfile(output_json):
        with open(output_json, "r", encoding="utf-8") as outf:
            output_list = json.load(outf)
        processed = set(r['key'] for r in output_list)
        logging.info(f"[{bucket}] Resuming: {len(processed)} files previously classified.")
    else:
        output_list = []
        processed = set()

    todo_records = [r for r in records if r['key'] not in processed]
    logging.info(f"[{bucket}] {len(todo_records)} files left to classify in parallel.")

    N = 12
    count = 0
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        futures = {executor.submit(process_one_file, bucket, rec, processed, compliance_matrix): rec for rec in todo_records}
        for future in as_completed(futures):
            outrec = future.result()
            if outrec:
                output_list.append(outrec)
                processed.add(outrec['key'])
                count += 1
                if count % N == 0:
                    with open(output_json, "w", encoding="utf-8") as outf:
                        json.dump(output_list, outf, indent=2)
                    logging.info(f"[{bucket}] {count} new this run, total {len(output_list)} files classified.")
    with open(output_json, "w", encoding="utf-8") as outf:
        json.dump(output_list, outf, indent=2)
    logging.info(f"[{bucket}] Completed. {len(output_list)} files classified.")