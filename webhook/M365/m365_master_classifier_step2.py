# Now your script works from either:

# /path/to/{config_id} → uses this as the config directory
# /path/to/somewhere → will look for any {config_id} folder below
# No more "not found" error if you cd to the correct config folder!

# Your error is because your script is looking for a subfolder named M365 inside /home/cybersecai/htdocs/www.cybersecai.io/webhook/M365, but you're already in the M365 folder itself. So it fails to find itself.

import os
import sys
import json
import logging
import io
import re
import time
import base64
import tempfile
import subprocess
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed

import requests
from requests.exceptions import RequestException

from docx import Document
from PyPDF2 import PdfReader

# Optional/extra filetype libraries
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None
try:
    from odf.opendocument import load as odf_load
    from odf.text import P
except ImportError:
    odf_load = None
    P = None
try:
    import extract_msg
except ImportError:
    extract_msg = None
try:
    from msg_parser import MsOxMessage as MsgParserMessage
except ImportError:
    MsgParserMessage = None
try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    import xlrd  # only supports .xls when xlrd<2.0
except ImportError:
    xlrd = None

import email
import email.policy

# OCR-related libs
try:
    from PIL import Image, ImageFile
    ImageFile.LOAD_TRUNCATED_IMAGES = True
except ImportError:
    Image = None
try:
    import pytesseract
except ImportError:
    pytesseract = None

# Optional PDF generation fallback (if LibreOffice not installed)
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import simpleSplit
except ImportError:
    canvas = None

# OpenAI (support both legacy and 1.x SDK)
try:
    from openai import OpenAI as OpenAIClientV1
except ImportError:
    OpenAIClientV1 = None
try:
    import openai as openai_legacy
except ImportError:
    openai_legacy = None

logging.basicConfig(format="%(asctime)s %(levelname)s %(message)s", level=logging.INFO)

ALLOWED_EXT = (
    '.docx', '.txt', '.pdf', '.csv', '.json', '.md',
    '.xls', '.xlsx', '.ppt', '.pptx', '.rtf', '.odt', '.ods', '.eml', '.msg',
    # images for OCR
    '.jpg', '.jpeg', '.png', '.tif', '.tiff', '.gif', '.bmp', '.webp'
)

MAX_OUT_CHARS = 15000
MAX_WORKERS = 8
OCR_MAX_BYTES = 10_000_000  # upper limit for OCR bytes we try to process locally
MAX_IMAGE_LONG_EDGE = int(os.getenv("OAI_MAX_IMAGE_LONG_EDGE", "2000"))  # px for OpenAI Vision downsizing
JPEG_QUALITY = int(os.getenv("OAI_JPEG_QUALITY", "85"))

# OpenAI model names
OAI_VISION_MODEL = os.getenv("OAI_VISION_MODEL", "gpt-4o-mini")
OAI_CLASSIFY_MODEL = os.getenv("OAI_CLASSIFY_MODEL", "gpt-4.1")

def abort(msg, code=1):
    logging.error(msg)
    sys.exit(code)

def which(bin_name):
    return shutil.which(bin_name)

HAS_SOFFICE = which("soffice") is not None
HAS_MSGCONVERT = which("msgconvert") is not None  # from libemail-outlook-message-perl (Linux)

def find_config_folder(config_id, root):
    for rootdir, dirs, files in os.walk(root):
        if os.path.basename(rootdir) == config_id:
            return os.path.abspath(rootdir)
    abort(f"Could not find config folder '{config_id}' in '{root}'")

def load_json(path, desc):
    if not os.path.isfile(path):
        abort(f"{desc} not found: {path}")
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

# ---------------------- OpenAI helpers ---------------------- #
def get_openai_client():
    api_key = os.getenv("OPENAI_API_KEY") or os.getenv("OPENAI_APIKEY") or os.getenv("OPENAI_KEY")
    if not api_key:
        return None, None
    if OpenAIClientV1 is not None:
        return OpenAIClientV1(api_key=api_key), "v1"
    if openai_legacy is not None:
        openai_legacy.api_key = api_key
        return openai_legacy, "legacy"
    return None, None

def prepare_image_for_vision(image_bytes):
    if Image is None:
        return image_bytes, "image/png"
    try:
        with Image.open(io.BytesIO(image_bytes)) as im:
            if im.mode not in ("RGB", "L"):
                im = im.convert("RGB")
            w, h = im.size
            long_edge = max(w, h)
            if long_edge > MAX_IMAGE_LONG_EDGE:
                scale = MAX_IMAGE_LONG_EDGE / float(long_edge)
                im = im.resize((max(1, int(w*scale)), max(1, int(h*scale))))
            buf = io.BytesIO()
            fmt = "JPEG" if im.mode in ("RGB", "L") else "PNG"
            if fmt == "JPEG":
                im.save(buf, format=fmt, quality=JPEG_QUALITY, optimize=True)
                mime = "image/jpeg"
            else:
                im.save(buf, format="PNG", optimize=True)
                mime = "image/png"
            return buf.getvalue(), mime
    except Exception as e:
        logging.warning(f"prepare_image_for_vision failed, using original: {e}")
        return image_bytes, "application/octet-stream"

def openai_vision_ocr(image_bytes):
    client, mode = get_openai_client()
    if client is None:
        return None
    try:
        prepped, mime = prepare_image_for_vision(image_bytes)
        data_url = f"data:{mime};base64,{base64.b64encode(prepped).decode('ascii')}"
        user_content = [
            {"type": "text", "text": "Extract all visible text from the image. Return only the raw text."},
            {"type": "image_url", "image_url": {"url": data_url}},
        ]
        if mode == "v1":
            resp = client.chat.completions.create(
                model=OAI_VISION_MODEL,
                messages=[{"role": "user", "content": user_content}],
                temperature=0,
            )
            return (resp.choices[0].message.content or "").strip()
        else:
            resp = client.ChatCompletion.create(
                model=OAI_VISION_MODEL,
                messages=[{"role": "user", "content": user_content}],
                temperature=0,
            )
            return (resp["choices"][0]["message"]["content"] or "").strip()
    except Exception as e:
        logging.warning(f"OpenAI Vision OCR error: {e}")
        return None

# OCR helpers
def pytesseract_ocr(image_bytes):
    if Image is None or pytesseract is None:
        return None
    try:
        with Image.open(io.BytesIO(image_bytes)) as im:
            return pytesseract.image_to_string(im)
    except Exception as e:
        logging.warning(f"pytesseract OCR failed: {e}")
        return None

def ocr_image_bytes(blob, ext):
    # Try OpenAI Vision first, then pytesseract
    try:
        text = openai_vision_ocr(blob)
        if text and text.strip():
            return text[:MAX_OUT_CHARS]
    except Exception as e:
        logging.warning(f"OpenAI Vision OCR fallback failed: {e}")
    text = pytesseract_ocr(blob)
    return text[:MAX_OUT_CHARS] if text else None

# ---------------------- Conversion helpers ---------------------- #
def libreoffice_convert(input_path, outdir, target):
    """
    Generic LO headless convert. target examples:
      - 'pdf:writer_pdf_Export' (Writer)
      - 'pdf:calc_pdf_Export' (Calc)
      - 'pdf:impress_pdf_Export' (Impress)
      - 'pdf' (auto)
      - 'xlsx' / 'docx' / 'csv' etc.
    """
    if not HAS_SOFFICE:
        return None
    cmd = ["soffice", "--headless", "--norestore", "--convert-to", target, "--outdir", outdir, input_path]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=240)
        base = os.path.splitext(os.path.basename(input_path))[0]
        ext = target.split(":")[0]
        if ext in ("pdf", "docx", "xlsx", "csv", "odt", "ods"):
            outpath = os.path.join(outdir, base + "." + ext)
            if os.path.isfile(outpath):
                return outpath
        for fn in os.listdir(outdir):
            if fn.startswith(base + "."):
                candidate = os.path.join(outdir, fn)
                if os.path.isfile(candidate):
                    return candidate
    except Exception as e:
        logging.warning(f"LibreOffice convert failed ({target}): {e}")
    return None

def libreoffice_convert_to_pdf(input_path, work_dir=None, kind_hint=None):
    outdir = work_dir or os.path.dirname(input_path)
    target = "pdf"
    if kind_hint == "writer":
        target = "pdf:writer_pdf_Export"
    elif kind_hint == "calc":
        target = "pdf:calc_pdf_Export"
    elif kind_hint == "impress":
        target = "pdf:impress_pdf_Export"
    return libreoffice_convert(input_path, outdir, target)

def extract_text_from_pdf_bytes(pdf_bytes, ocr_if_empty=True):
    text = ""
    try:
        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            text += (page.extract_text() or "") + "\n"
    except Exception as e:
        logging.warning(f"PDF parse error: {e}")
    # Note: Additional OCR of PDF pages (pdf2image) can be added if needed.
    return text

# ---------------------- .msg/.eml helpers ---------------------- #
def parse_msg_with_extract_msg(path):
    if extract_msg is None:
        return None
    try:
        msg = extract_msg.Message(path)
        subject = msg.subject or ""
        sender = getattr(msg, "sender", "") or getattr(msg, "sender_email", "")
        to = msg.to or ""
        date = msg.date or ""
        body = msg.body or ""
        if not body:
            body = getattr(msg, "htmlBody", "") or ""
        if not body and hasattr(msg, "rtfBody") and msg.rtfBody:
            try:
                if isinstance(msg.rtfBody, bytes):
                    body = rtf_to_text(msg.rtfBody.decode(errors="ignore")) if rtf_to_text else ""
                else:
                    body = rtf_to_text(msg.rtfBody) if rtf_to_text else ""
            except Exception:
                pass
        text = f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nDate: {date}\n\n{body}"
        try:
            msg.close()
        except Exception:
            pass
        return text
    except Exception as e:
        logging.warning(f"extract_msg parse error: {e}")
        return None

def parse_msg_with_msg_parser(path):
    if MsgParserMessage is None:
        return None
    try:
        m = MsgParserMessage(path)
        m_message = m.get_email_data()
        subject = m_message.get("subject", "") or ""
        sender = m_message.get("sender", "") or ""
        to = ", ".join(m_message.get("to", []) or [])
        date = m_message.get("date", "") or ""
        body = (m_message.get("body", "") or "") + "\n" + (m_message.get("body_html", "") or "")
        text = f"Subject: {subject}\nFrom: {sender}\nTo: {to}\nDate: {date}\n\n{body}"
        return text
    except Exception as e:
        logging.warning(f"msg_parser parse error: {e}")
        return None

def convert_msg_to_eml_with_msgconvert(msg_path, out_dir):
    if not HAS_MSGCONVERT:
        return None
    try:
        res = subprocess.run(["msgconvert", msg_path], cwd=out_dir, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)
        if res.returncode != 0:
            logging.warning(f"msgconvert failed: {res.stderr.decode(errors='ignore')}")
            return None
        base = os.path.splitext(os.path.basename(msg_path))[0]
        for fn in os.listdir(out_dir):
            if fn.lower().startswith(base.lower()) and fn.lower().endswith(".eml"):
                return os.path.join(out_dir, fn)
        emls = [os.path.join(out_dir, fn) for fn in os.listdir(out_dir) if fn.lower().endswith(".eml")]
        if emls:
            return sorted(emls, key=lambda p: os.path.getmtime(p), reverse=True)[0]
    except Exception as e:
        logging.warning(f"msgconvert exception: {e}")
    return None

def parse_eml_bytes(content):
    try:
        msg = email.message_from_bytes(content, policy=email.policy.default)
        text_parts = []
        if msg['Subject']:
            text_parts.append("Subject: " + str(msg['Subject']))
        if msg['From']:
            text_parts.append("From: " + str(msg['From']))
        if msg['To']:
            text_parts.append("To: " + str(msg['To']))
        if msg['Date']:
            text_parts.append("Date: " + str(msg['Date']))
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                disp = (part.get("Content-Disposition") or "").lower()
                if ctype == "text/plain" and "attachment" not in disp:
                    try:
                        body += part.get_content()
                    except Exception:
                        try:
                            body += (part.get_payload(decode=True) or b"").decode(errors="ignore")
                        except Exception:
                            pass
        else:
            try:
                body = msg.get_content()
            except Exception:
                body = (msg.get_payload(decode=True) or b"").decode(errors="ignore")
        text_parts.append(str(body))
        return "\n".join(text_parts)
    except Exception as e:
        logging.warning(f"EML parse error: {e}")
        return None

def simple_extract_msg_text(path):
    text = parse_msg_with_extract_msg(path)
    if text and text.strip():
        return text[:MAX_OUT_CHARS]
    text = parse_msg_with_msg_parser(path)
    if text and text.strip():
        return text[:MAX_OUT_CHARS]
    with tempfile.TemporaryDirectory() as td:
        eml_path = convert_msg_to_eml_with_msgconvert(path, td)
        if eml_path and os.path.isfile(eml_path):
            try:
                with open(eml_path, "rb") as f:
                    content = f.read()
                text = parse_eml_bytes(content)
                if text and text.strip():
                    return text[:MAX_OUT_CHARS]
            except Exception as e:
                logging.warning(f"EML read after msgconvert failed: {e}")
    return None

def html_escape(s):
    try:
        import html
        return html.escape(s)
    except Exception:
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def write_text_as_html_file(text, out_dir, base="email_export"):
    html_content = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>Email Export</title>
<style>
body {{ font-family: Arial, sans-serif; white-space: pre-wrap; }}
.header {{ font-weight: bold; margin-bottom: 10px; }}
</style>
</head>
<body>
{html_escape(text)}
</body>
</html>
"""
    html_path = os.path.join(out_dir, f"{base}.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    return html_path

def reportlab_write_pdf(text, out_pdf):
    if canvas is None:
        return False
    try:
        c = canvas.Canvas(out_pdf, pagesize=letter)
        width, height = letter
        margin = 36  # half inch
        max_width = width - 2*margin
        y = height - margin
        c.setFont("Helvetica", 10)
        lines = text.splitlines() or [text]
        for line in lines:
            wrapped = simpleSplit(line, "Helvetica", 10, max_width)
            for wline in wrapped:
                if y <= margin:
                    c.showPage()
                    c.setFont("Helvetica", 10)
                    y = height - margin
                c.drawString(margin, y, wline)
                y -= 12
        c.save()
        return True
    except Exception as e:
        logging.warning(f"ReportLab PDF generation failed: {e}")
        return False

def convert_msg_to_pdf(msg_path, work_dir=None):
    td_created = False
    if work_dir is None:
        work_dir = tempfile.mkdtemp()
        td_created = True
    pdf_path = None
    try:
        eml_path = convert_msg_to_eml_with_msgconvert(msg_path, work_dir)
        if eml_path and os.path.isfile(eml_path):
            pdf_path = libreoffice_convert_to_pdf(eml_path, work_dir=work_dir, kind_hint="writer")
            if pdf_path and os.path.isfile(pdf_path):
                return pdf_path

        pdf_path = libreoffice_convert_to_pdf(msg_path, work_dir=work_dir, kind_hint="writer")
        if pdf_path and os.path.isfile(pdf_path):
            return pdf_path

        text = simple_extract_msg_text(msg_path) or ""
        if text.strip():
            html_path = write_text_as_html_file(text, work_dir, base="email_from_msg")
            pdf_path = libreoffice_convert_to_pdf(html_path, work_dir=work_dir, kind_hint="writer")
            if pdf_path and os.path.isfile(pdf_path):
                return pdf_path
            rl_pdf = os.path.join(work_dir, "email_from_msg_reportlab.pdf")
            if reportlab_write_pdf(text, rl_pdf):
                return rl_pdf

        logging.warning("convert_msg_to_pdf: All PDF conversion paths failed.")
        return None
    finally:
        if td_created:
            if pdf_path is None:
                try:
                    for fn in os.listdir(work_dir):
                        os.unlink(os.path.join(work_dir, fn))
                    os.rmdir(work_dir)
                except Exception:
                    pass

# ---------------------- Excel helpers ---------------------- #
def excel_to_text_with_openpyxl(xlsx_path, max_rows=1000):
    if openpyxl is None:
        return None
    try:
        wb = openpyxl.load_workbook(xlsx_path, read_only=True, data_only=True)
        pieces = []
        for ws in wb.worksheets:
            rows = []
            rcount = 0
            for row in ws.iter_rows(values_only=True):
                rcount += 1
                if rcount > max_rows:
                    break
                vals = [(str(v) if v is not None else "") for v in row]
                rows.append(",".join(vals))
            if rows:
                pieces.append(f"=== Sheet: {ws.title} ===\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(pieces)
    except Exception as e:
        logging.warning(f"openpyxl extract error: {e}")
        return None

def excel_xls_to_text_with_xlrd(xls_path, max_rows=1000):
    if xlrd is None:
        return None
    try:
        book = xlrd.open_workbook(xls_path, on_demand=True)
        pieces = []
        for sheet_name in book.sheet_names():
            sh = book.sheet_by_name(sheet_name)
            rows = []
            for r in range(min(sh.nrows, max_rows)):
                vals = []
                for c in range(sh.ncols):
                    cell = sh.cell(r, c)
                    if cell.ctype == xlrd.XL_CELL_DATE:
                        try:
                            from datetime import datetime
                            dt = xlrd.xldate_as_tuple(cell.value, book.datemode)
                            vals.append(str(datetime(*dt)))
                        except Exception:
                            vals.append(str(cell.value))
                    else:
                        vals.append(str(cell.value) if cell.value is not None else "")
                rows.append(",".join(vals))
            if rows:
                pieces.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
        book.release_resources()
        return "\n\n".join(pieces)
    except Exception as e:
        logging.warning(f"xlrd extract error: {e}")
        return None

def excel_bytes_to_text(blob, ext):
    with tempfile.TemporaryDirectory() as td:
        in_path = os.path.join(td, f"input{ext}")
        with open(in_path, "wb") as f:
            f.write(blob)

        if ext == ".xlsx":
            txt = excel_to_text_with_openpyxl(in_path)
            if txt and txt.strip():
                return txt

            if pd is not None:
                try:
                    xls = pd.ExcelFile(in_path, engine='openpyxl')
                    pieces = []
                    for sheet in xls.sheet_names:
                        try:
                            df = pd.read_excel(xls, sheet_name=sheet, dtype=str, nrows=1000, na_filter=False)
                            if not df.empty:
                                csv_text = df.to_csv(index=False)
                                if csv_text.strip():
                                    pieces.append(f"=== Sheet: {sheet} ===\n{csv_text}")
                        except Exception as se:
                            logging.warning(f"Pandas Excel sheet parse error ({sheet}): {se}")
                    if pieces:
                        return "\n\n".join(pieces)
                except Exception as e:
                    logging.warning(f"Pandas Excel xlsx parse error: {e}")

        else:  # .xls
            txt = excel_xls_to_text_with_xlrd(in_path)
            if txt and txt.strip():
                return txt

            if pd is not None:
                try:
                    xls = pd.ExcelFile(in_path)
                    pieces = []
                    for sheet in xls.sheet_names:
                        try:
                            df = pd.read_excel(xls, sheet_name=sheet, dtype=str, nrows=1000, na_filter=False)
                            if not df.empty:
                                csv_text = df.to_csv(index=False)
                                if csv_text.strip():
                                    pieces.append(f"=== Sheet: {sheet} ===\n{csv_text}")
                        except Exception as se:
                            logging.warning(f"Pandas Excel sheet parse error ({sheet}): {se}")
                    if pieces:
                        return "\n\n".join(pieces)
                except Exception as e:
                    logging.warning(f"Pandas Excel xls parse error: {e}")

        if HAS_SOFFICE:
            xlsx_path = libreoffice_convert(in_path, td, "xlsx")
            if xlsx_path and os.path.isfile(xlsx_path):
                txt = excel_to_text_with_openpyxl(xlsx_path)
                if txt and txt.strip():
                    return txt
            pdf_path = libreoffice_convert_to_pdf(in_path, work_dir=td, kind_hint="calc")
            if pdf_path and os.path.isfile(pdf_path):
                with open(pdf_path, "rb") as pf:
                    pdf_bytes = pf.read()
                text = extract_text_from_pdf_bytes(pdf_bytes)
                if text and text.strip():
                    return text

    return None

# ---------------------- PPT helpers ---------------------- #
def parse_pptx_from_file(path):
    if Presentation is not None:
        try:
            prs = Presentation(path)
            lines = []
            for slide in prs.slides:
                if getattr(slide.shapes, "title", None) and getattr(slide.shapes.title, "text", None):
                    lines.append(slide.shapes.title.text)
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and getattr(shape, "text", None):
                        lines.append(shape.text)
                    if hasattr(shape, "table") and shape.table is not None:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text:
                                    lines.append(cell.text)
            combined = "\n".join(lines).strip()
            if combined:
                return combined[:MAX_OUT_CHARS]
        except Exception as e:
            logging.warning(f"python-pptx parse error: {e}")
    pdf_path = libreoffice_convert_to_pdf(path, work_dir=os.path.dirname(path), kind_hint="impress")
    if pdf_path and os.path.isfile(pdf_path):
        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()
        text = extract_text_from_pdf_bytes(pdf_bytes)
        try:
            os.unlink(pdf_path)
        except Exception:
            pass
        return text[:MAX_OUT_CHARS] if text else None
    return None

def parse_ppt_from_bytes(blob, ext):
    with tempfile.TemporaryDirectory() as td:
        in_path = os.path.join(td, f"input{ext}")
        with open(in_path, "wb") as f:
            f.write(blob)
        pdf_path = libreoffice_convert_to_pdf(in_path, work_dir=td, kind_hint="impress")
        if pdf_path and os.path.isfile(pdf_path):
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()
            text = extract_text_from_pdf_bytes(pdf_bytes)
            return text[:MAX_OUT_CHARS] if text else None
    return None

# ---------------------- Microsoft Graph Client ---------------------- #
class GraphClient:
    def __init__(self, tenant_id, client_id, client_secret, timeout=60):
        self.tenant_id = tenant_id
        self.client_id = client_id
        self.client_secret = client_secret
        self.token = None
        self.session = requests.Session()
        self.timeout = timeout
        self._get_access_token()

    def _get_access_token(self):
        url = f"https://login.microsoftonline.com/{self.tenant_id}/oauth2/v2.0/token"
        data = {
            'grant_type': 'client_credentials',
            'client_id': self.client_id,
            'client_secret': self.client_secret,
            'scope': 'https://graph.microsoft.com/.default'
        }
        try:
            resp = self.session.post(url, data=data, timeout=20)
            resp.raise_for_status()
            self.token = resp.json()['access_token']
            logging.info("Obtained Microsoft Graph access token.")
        except Exception as e:
            abort(f"Failed to obtain access token: {e}")

    def _request(self, method, url, retries=4):
        last_exc = None
        for attempt in range(retries):
            try:
                headers = {"Authorization": f"Bearer {self.token}"}
                resp = self.session.request(method, url, headers=headers, timeout=self.timeout, allow_redirects=True)
                if resp.status_code == 401:
                    logging.warning("Graph 401 Unauthorized, refreshing token...")
                    self._get_access_token()
                    continue
                if resp.status_code in (429, 500, 502, 503, 504):
                    retry_after = int(resp.headers.get("Retry-After", "0"))
                    delay = retry_after if retry_after > 0 else (1.5 ** attempt)
                    logging.warning(f"Graph {resp.status_code} for {url}. Retrying in {delay:.1f}s")
                    time.sleep(delay)
                    continue
                resp.raise_for_status()
                return resp
            except RequestException as e:
                last_exc = e
                delay = 1.5 ** attempt
                logging.warning(f"Graph request error: {e}. Retrying in {delay:.1f}s")
                time.sleep(delay)
        if last_exc:
            raise last_exc
        raise Exception("Graph request failed after retries")

    def download_file_bytes(self, file_info):
        """
        file_info can represent OneDrive user drive item, or SharePoint site drive item.
        Supported combinations:
          - user_id + file_id
          - site_id + drive_id + file_id
          - drive_id + file_id
        """
        url = None
        if file_info.get('user_id') and file_info.get('file_id'):
            url = f"https://graph.microsoft.com/v1.0/users/{file_info['user_id']}/drive/items/{file_info['file_id']}/content"
        elif file_info.get('site_id') and file_info.get('drive_id') and file_info.get('file_id'):
            url = f"https://graph.microsoft.com/v1.0/sites/{file_info['site_id']}/drives/{file_info['drive_id']}/items/{file_info['file_id']}/content"
        elif file_info.get('drive_id') and file_info.get('file_id'):
            url = f"https://graph.microsoft.com/v1.0/drives/{file_info['drive_id']}/items/{file_info['file_id']}/content"
        else:
            raise ValueError("Record missing user_id or site_id/drive_id for file")
        resp = self._request("GET", url)
        return resp.content

# ---------------------- File content extraction ---------------------- #
def fetch_file_content(graph_client, record):
    """
    Download content via Graph and extract text based on file type.
    """
    fname = record.get('file_name', '')
    ext = os.path.splitext(fname)[1].lower()
    if ext not in ALLOWED_EXT or record.get('file_type', '').lower() == 'folder':
        return None
    try:
        content = graph_client.download_file_bytes(record)

        # Images -> OCR
        if ext in ('.jpg', '.jpeg', '.png', '.tif', '.tiff', '.gif', '.bmp', '.webp'):
            blob = content[:OCR_MAX_BYTES]
            return ocr_image_bytes(blob, ext)

        # Plain text family
        if ext in ('.txt', '.csv', '.json', '.md'):
            try:
                return content.decode('utf-8')[:MAX_OUT_CHARS]
            except Exception:
                return content.decode('latin-1', errors="replace")[:MAX_OUT_CHARS]

        elif ext == '.docx':
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmpf:
                tmpf.write(content)
                tmpname = tmpf.name
            try:
                doc = Document(tmpname)
                return "\n".join(p.text for p in doc.paragraphs)[:MAX_OUT_CHARS]
            finally:
                try: os.unlink(tmpname)
                except Exception: pass

        elif ext == '.pdf':
            text = extract_text_from_pdf_bytes(content)
            if len((text or "").strip()) < 80:
                # as scanned fallback, try OCR via Vision + pytesseract on rendered image bytes if needed
                # here we simply try Vision directly on the PDF bytes (may or may not work),
                # better approach: render pages into images. For now, use Vision on bytes:
                ocr_txt = openai_vision_ocr(content)
                if ocr_txt and len(ocr_txt.strip()) > len((text or "").strip()):
                    text = ocr_txt
            return text[:MAX_OUT_CHARS] if text else None

        elif ext in ('.xls', '.xlsx'):
            text = excel_bytes_to_text(content, ext)
            return text[:MAX_OUT_CHARS] if text else None

        elif ext in ('.pptx', '.ppt'):
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmpf:
                tmpf.write(content)
                inpath = tmpf.name
            try:
                if ext == '.pptx':
                    text = parse_pptx_from_file(inpath)
                else:
                    text = parse_ppt_from_bytes(content, ext)
                if (not text) or len(text.strip()) < 10:
                    pdf_path = libreoffice_convert_to_pdf(inpath, kind_hint="impress")
                    if pdf_path and os.path.isfile(pdf_path):
                        with open(pdf_path, "rb") as pf:
                            pdf_bytes = pf.read()
                        text = extract_text_from_pdf_bytes(pdf_bytes)
                return text[:MAX_OUT_CHARS] if text else None
            finally:
                try: os.unlink(inpath)
                except Exception: pass

        elif ext == '.rtf' and rtf_to_text is not None:
            try:
                return rtf_to_text(content.decode(errors="ignore"))[:MAX_OUT_CHARS]
            except Exception as e:
                logging.warning(f"RTF parse error: {e}")
                return None

        elif ext in ('.odt', '.ods') and odf_load is not None and P is not None:
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmpf:
                tmpf.write(content)
                tmpname = tmpf.name
            try:
                doc = odf_load(tmpname)
                paragraphs = []
                for elem in doc.getElementsByType(P):
                    try:
                        paragraphs.append(elem.firstChild.data if elem.firstChild else "")
                    except Exception:
                        paragraphs.append(str(elem))
                combined = "\n".join(paragraphs)
                return combined[:MAX_OUT_CHARS] if combined else None
            except Exception as e:
                logging.warning(f"ODF parse error: {e}")
                return None
            finally:
                try: os.unlink(tmpname)
                except Exception: pass

        elif ext == '.eml':
            text = parse_eml_bytes(content)
            return text[:MAX_OUT_CHARS] if text else None

        elif ext == '.msg':
            with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as tmpf:
                tmpf.write(content)
                tmpname = tmpf.name
            try:
                pdf_path = convert_msg_to_pdf(tmpname)
                text = None
                if pdf_path and os.path.isfile(pdf_path):
                    try:
                        with open(pdf_path, "rb") as pf:
                            pdf_bytes = pf.read()
                        text = extract_text_from_pdf_bytes(pdf_bytes)
                        if (not text) or len(text.strip()) < 50:
                            extra_text = simple_extract_msg_text(tmpname)
                            if extra_text:
                                text = (text or "") + "\n\n" + extra_text
                    except Exception as e:
                        logging.warning(f"Reading/extracting from generated PDF failed: {e}")
                else:
                    logging.warning("MSG->PDF conversion failed; falling back to direct text extraction.")
                    text = simple_extract_msg_text(tmpname)
                return text[:MAX_OUT_CHARS] if text else None
            finally:
                try: os.unlink(tmpname)
                except Exception: pass

        else:
            return ""
    except Exception as e:
        logging.warning(f"Failed to extract content for {fname}: {e}")
        return None

# ---------------------- Classification prompt ---------------------- #
def build_prompt(compliance_matrix):
    s = (
        "You are an expert compliance auditor and cyber expert. Based on your expert knowledge and using the table below for standards and regulated data fields as a reference,\n"
        "analyze the provided document (or data) for regulated data types and regulated data fields. When reviewing, pay close attention to the following and document your reasoning:\n"
        "- If the file contains or appears to contain only training material, instructional content with fake/sample/demo data, or random values (not real production data), clearly indicate this in your results and reasoning.\n"
        "- If regulated fields are present but there are no real (production or sensitive) values associated, rate the risk and data classification LOWER than if real production data is present.\n"
        "- If the data is explicitly marked as sample data, test records, or is clearly intended for demonstration or training, reduce the file's risk and sensitivity rating accordingly and document this in your explanation.\n"
        "- Only assign High or Medium risk and 'sensitive' classifications where there is strong evidence of actual regulated data populated in those fields.\n"
        "\n"
        "For every compliance standard with a match, create a results array entry containing:\n"
        "  - \"standard\": name of the standard\n"
        "  - \"jurisdiction\": the law's jurisdiction/region\n"
        "  - \"detected_fields\": list of regulated field names found in this document\n"
        "  - \"risk\": High, Medium, Low, or None, based on the actual data detected and the nature of that data (with sample/demo/fake data as Low or None)\n"
        "After reviewing all, output the following as ordered fields at the END of the JSON object:\n"
        "  - \"auditor_agent_view\": In 2-3 sentences, explain your reasoning (plain English, for auditors), specifically highlight if the content appears to be training/sample/demo, or fake/random data.\n"
        "  - \"likely_data_subject_area\": State data subject area, e.g., customer data, financial records, health information, etc.\n"
        "  - \"data_classification\": Classify Data, e.g., Highly Sensitive, Confidential, Internal Use, Public etc. Use a lower classification if the file only contains training/sample/demo data or regulated fields without actual sensitive values.\n"
        "  - \"overall_risk_rating\": string, High, Medium, Low, or None (the highest risk found overall; use Low or None if only sample/demo/empty data)\n"
        "  - \"hacker_interest\": State How Hacker or bad actors will potentially exploit this information if they get access to it. If the file only has demo/sample data, state 'Little to no value to hackers.'\n"
        "  - \"cyber_proposed_controls\": List proposed cyber controls leveraging ISO standard and NIST, based on data classification and risk rating as a cyber expert\n"
        "  - \"auditor_proposed_action\": string, short action phrase for remediation based on auditor expert knowledge, e.g., \"notify owner\", \"quarantine file\", \"no action required\" etc. For training/sample/demo data, recommend \"no action required\" or similar.\n"
        "Respond with a SINGLE pure JSON object only, in this format:\n"
        "{\n"
        "  \"results\": [\n"
        "     {\"standard\":\"\",\"jurisdiction\":\"\",\"detected_fields\":[],\"risk\":\"\",\"auditor_agent_view\":\"\"}, ...\n"
        "  ],\n"
        "  \"auditor_agent_view\": \"\",\n"
        "  \"likely_data_subject_area\": \"\",\n"
        "  \"data_classification\": \"\",\n"
        "  \"overall_risk_rating\": \"\",\n"
        "  \"hacker_interest\": \"\",\n"
        "  \"cyber_proposed_controls\": \"\",\n"
        "  \"auditor_proposed_action\": \"\"\n"
        "}\n"
        "Here is your standards/fields matrix:\n"
    )
    for entry in compliance_matrix:
        s += (
            f"Standard: {entry['standard']} | Jurisdiction: {entry['jurisdiction']} | "
            f"Fields: {', '.join(entry['fields'])}\n"
        )
    return s

def gpt_classify_file(text, compliance_matrix):
    client, mode = get_openai_client()
    if client is None:
        logging.error("OpenAI SDK not installed or OPENAI_API_KEY not set; cannot classify.")
        return "OpenAI not available"
    system_prompt = build_prompt(compliance_matrix)
    user = f"Classify the following file content for regulated data fields (see system context):\n\n{text[:MAX_OUT_CHARS]}"
    try:
        if mode == "v1":
            response = client.chat.completions.create(
                model=OAI_CLASSIFY_MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user}
                ],
                temperature=0.0,
            )
            return (response.choices[0].message.content or "").strip()
        else:
            response = client.ChatCompletion.create(
                model=OAI_CLASSIFY_MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user}
                ],
                temperature=0.0,
            )
            return (response["choices"][0]["message"]["content"] or "").strip()
    except Exception as e:
        logging.error(f"GPT API error: {e}")
        return f"GPT API error: {e}"

# ---------------------- Orchestration ---------------------- #
def which_m365_jsons(config_dir, config_id):
    # Matches any prefix, then _m365_files_list_<CONFIG_ID>.json
    pat = re.compile(r"(.+)_m365_files_list_" + re.escape(config_id) + r"\.json$")
    return [fn for fn in os.listdir(config_dir) if pat.match(fn)], pat

def process_one_file(graph_client, record, compliance_matrix):
    text = fetch_file_content(graph_client, record)
    if text and text.strip():
        try:
            res = gpt_classify_file(text[:MAX_OUT_CHARS], compliance_matrix)
            outrec = dict(record)
            outrec['llm_response'] = res
            logging.info(f"Classified: {record.get('file_name')}")
            return outrec
        except Exception as e:
            logging.error(f"Failed GPT for {record.get('file_name')}: {e}")
            outrec = dict(record)
            outrec['llm_response'] = f"Error: {e}"
            return outrec
    else:
        logging.info(f"No usable text for file: {record.get('file_name')}")
        return None

def main():
    if len(sys.argv) < 2:
        abort("Usage: python3 m365_content_extract_classify.py <CONFIG_ID> [<search_root>]")

    config_id = sys.argv[1]
    search_root = sys.argv[2] if len(sys.argv) > 2 else '/home/cybersecai/htdocs/www.cybersecai.io/webhook/M365'

    config_dir = find_config_folder(config_id, search_root)
    graph_dir  = os.path.join(config_dir, "graph")
    os.makedirs(graph_dir, exist_ok=True)

    secrets_json = os.path.join(config_dir, f"{config_id}_secrets.json")
    secrets = load_json(secrets_json, "M365 Secrets file")
    for key in ("TENANT_ID", "CLIENT_ID", "CLIENT_SECRET"):
        if key not in secrets or not secrets[key]:
            abort(f"Missing or blank: {key} in secrets file")

    compliance_matrix = load_json(os.path.join(config_dir, "compliance_matrix.json"), "Compliance matrix")

    if not (os.getenv("OPENAI_API_KEY") or os.getenv("OPENAI_APIKEY") or os.getenv("OPENAI_KEY")):
        abort("Please set your OpenAI API key in OPENAI_API_KEY (or OPENAI_APIKEY/OPENAI_KEY) environment variable.")

    graph_client = GraphClient(secrets['TENANT_ID'], secrets['CLIENT_ID'], secrets['CLIENT_SECRET'])

    all_jsons, pat = which_m365_jsons(config_dir, config_id)
    if not all_jsons:
        logging.warning(f"No M365 file list JSONs found for config {config_id} in {config_dir}")
    else:
        logging.info(f"Found {len(all_jsons)} M365 file lists: {all_jsons}")

    for filelist_json in all_jsons:
        abs_filelist = os.path.join(config_dir, filelist_json)
        if not os.path.isfile(abs_filelist) or os.path.getsize(abs_filelist) == 0:
            logging.warning(f"Skipping empty or missing file list: {abs_filelist}")
            continue

        output_json = os.path.join(graph_dir, f"output_{filelist_json}")
        try:
            records = load_json(abs_filelist, f"M365 file list {filelist_json}")
            if not isinstance(records, list):
                logging.warning(f"File list {filelist_json} not a list; skipping.")
                continue
        except Exception as e:
            logging.error(f"Failed to read {abs_filelist}: {e}")
            continue

        # Load previous results to support delta by file_id + last_modified
        prev_by_id = {}
        if os.path.isfile(output_json):
            try:
                with open(output_json, "r", encoding="utf-8") as outf:
                    prev_list = json.load(outf)
                for rec in prev_list:
                    if 'file_id' in rec:
                        prev_by_id[rec['file_id']] = rec
                logging.info(f"Resuming: {len(prev_by_id)} files previously classified.")
            except Exception as e:
                logging.warning(f"Could not read previous output {output_json}: {e}")

        # Build todo and reuse lists
        todo_records = []
        reused_records = []
        present_ids = set()

        for rec in records:
            file_id = rec.get('file_id')
            if not file_id:
                continue
            present_ids.add(file_id)
            ext = os.path.splitext(rec.get('file_name', ''))[1].lower()
            if ext not in ALLOWED_EXT or rec.get('file_type','').lower() == 'folder':
                continue
            prev = prev_by_id.get(file_id)
            if prev and prev.get('last_modified') == rec.get('last_modified'):
                reused_records.append(prev)
            else:
                todo_records.append(rec)

        logging.info(f"[{filelist_json}] {len(todo_records)} files to classify, {len(reused_records)} reused.")

        output_list = list(reused_records)
        N = 12
        count_new = 0
        with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            futures = {executor.submit(process_one_file, graph_client, rec, compliance_matrix): rec for rec in todo_records}
            for future in as_completed(futures):
                outrec = future.result()
                if outrec:
                    output_list.append(outrec)
                    count_new += 1
                    if count_new % N == 0:
                        try:
                            with open(output_json, "w", encoding="utf-8") as outf:
                                json.dump(output_list, outf, indent=2)
                            logging.info(f"[{filelist_json}] {count_new} new this run, total {len(output_list)} files classified/reused.")
                        except Exception as e:
                            logging.error(f"Failed writing interim output: {e}")

        # Final write
        try:
            with open(output_json, "w", encoding="utf-8") as outf:
                json.dump(output_list, outf, indent=2)
            logging.info(f"[{filelist_json}] Completed. {len(output_list)} files classified/reused. Output: {output_json}")
        except Exception as e:
            logging.error(f"Failed writing final output {output_json}: {e}")

if __name__ == "__main__":
    main()