# Now your script works from either:

# /path/to/{config_id} → uses this as the config directory
# /path/to/somewhere → will look for any {config_id} folder below
# No more "not found" error if you cd to the correct config folder!

# Your error is because your script is looking for a subfolder named M365 inside /home/cybersecai/htdocs/www.cybersecai.io/webhook/M365, but you're already in the M365 folder itself. So it fails to find itself.

import os
import sys
import json
import requests
import re
import tempfile
import io
import logging
from datetime import datetime
from docx import Document
from PyPDF2 import PdfReader
import openai

# ==============================
# LOGGING SETUP & CONFIG
# ==============================
logging.basicConfig(
    format="%(asctime)s %(levelname)s %(message)s",
    level=logging.INFO,
)


ALLOWED_EXT = (
    '.docx', '.txt', '.pdf', '.csv', '.json', '.md',
    '.xls', '.xlsx', '.pptx', '.rtf', '.odt', '.ods', '.eml', '.msg'
)

# ---- NEW IMPORTS ----
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None
try:
    from odf.opendocument import load as odf_load
    from odf.text import P
except ImportError:
    odf_load = None
try:
    import extract_msg
except ImportError:
    extract_msg = None
import email
import email.policy
# ---------------------

MAX_OUT_CHARS = 15000
CHUNK_SIZE = 10000
OPENAI_MODEL = "gpt-4.1"

def abort(msg, code=1):
    logging.error(msg)
    raise SystemExit(code)

def find_config_folder(config_id, search_root):
    """
    Recursively search for /webhook/M365/<CONFIG_ID> under search_root.
    Returns the absolute path to the config directory.
    """
    logging.info(f"Searching for config folder '{config_id}' under '{search_root}'...")
    for root, dirs, files in os.walk(search_root):
        if os.path.basename(root) == config_id:
            logging.info(f"Located config directory: {root}")
            return os.path.abspath(root)
    abort(f"Config folder named '{config_id}' not found under '{search_root}'.")

def require_file(path, description):
    """Check that a file exists and log error if not."""
    if not os.path.isfile(path):
        abort(f"{description} not found: {path}")
    return path

def read_secrets(config_dir, config_id):
    secrets_path = require_file(os.path.join(config_dir, f"{config_id}_secrets.json"), "Secrets file")
    try:
        with open(secrets_path, 'r', encoding='utf-8') as f:
            d = json.load(f)
        for key in ("TENANT_ID", "CLIENT_ID", "CLIENT_SECRET"):
            assert key in d and d[key], f"Missing or blank: {key}"
        logging.info(f"Loaded secret credentials from {secrets_path}")
        return d
    except Exception as e:
        abort(f"Failed to load credentials from {secrets_path}: {e}")

def load_compliance_matrix(config_dir):
    cm_path = require_file(os.path.join(config_dir, "compliance_matrix.json"), "Compliance matrix")
    try:
        with open(cm_path, "r", encoding="utf-8") as f:
            matrix = json.load(f)
        assert isinstance(matrix, list) and all(isinstance(x, dict) for x in matrix)
        logging.info(f"Loaded compliance_matrix.json ({len(matrix)} standards)")
        return matrix
    except Exception as e:
        abort(f"Failed to load compliance_matrix.json: {e}")

def sanitize_filename(name):
    return re.sub(r'[^A-Za-z0-9._-]', '_', name)

def get_access_token(tenant_id, client_id, client_secret):
    url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
    data = {
        'grant_type': 'client_credentials',
        'client_id': client_id,
        'client_secret': client_secret,
        'scope': 'https://graph.microsoft.com/.default'
    }
    try:
        resp = requests.post(url, data=data, timeout=15)
        resp.raise_for_status()
        token = resp.json()['access_token']
        logging.info("Obtained Microsoft Graph access token.")
        return token
    except Exception as e:
        abort(f"Failed to obtain access token: {e}")




def fetch_file_content(file_info, access_token):
    ext = os.path.splitext(file_info['file_name'])[1].lower()
    base_url = ""
    def dual_log(msg, level='info'):
        from datetime import datetime
        LARAVEL_LOG = "/home/cybersecai/htdocs/www.cybersecai.io/storage/logs/laravel.log"
        timestr = datetime.utcnow().isoformat()
        text = f"[PYFETCH {timestr}] [{level.upper()}] {msg}"
        try:
            with open(LARAVEL_LOG, "a") as f:
                f.write(text + "\n")
        except Exception:
            pass
        getattr(logging, level.lower(), logging.info)(text)

    if 'user_id' in file_info and file_info.get('user_id') and 'file_id' in file_info:
        base_url = f"https://graph.microsoft.com/v1.0/users/{file_info['user_id']}/drive/items/{file_info['file_id']}/content"
        dual_log(f"Fetching content (OneDrive) for user {file_info['user_id']} file {file_info['file_id']}")
    elif 'site_id' in file_info and 'drive_id' in file_info and 'file_id' in file_info:
        base_url = f"https://graph.microsoft.com/v1.0/sites/{file_info['site_id']}/drives/{file_info['drive_id']}/items/{file_info['file_id']}/content"
        dual_log(f"Fetching content (SharePoint) for site {file_info['site_id']} drive {file_info['drive_id']} file {file_info['file_id']}")
    else:
        dual_log(f"Record missing user_id or site_id/drive_id for file {file_info.get('file_name','?')}, skipping.", "error")
        logging.warning("File record missing user_id or site_id/drive_id: %r", file_info)
        return ""

    headers = {"Authorization": f"Bearer {access_token}"}
    try:
        resp = requests.get(base_url, headers=headers)
        resp.raise_for_status()
        content = resp.content
    except Exception as e:
        dual_log(f"Failed to fetch content for file {file_info.get('file_name','?')}: {e}", "error")
        logging.error(f"Failed to fetch content: {e}")
        return ""

    # --- EXTENDED FILE SUPPORT ---
    if ext in ('.txt', '.csv', '.json', '.md'):
        try:
            return content.decode('utf-8')
        except Exception:
            return content.decode('latin-1', errors="replace")

    elif ext == '.docx':
        import tempfile
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
        try:
            tmp.write(content)
            tmp.close()
            doc = Document(tmp.name)
            text = "\n".join(p.text for p in doc.paragraphs)
            return text
        finally:
            os.unlink(tmp.name)

    elif ext == '.pdf':
        pdf_bytes = io.BytesIO(content)
        reader = PdfReader(pdf_bytes)
        text = ""
        for page in reader.pages:
            try:
                text += (page.extract_text() or "") + "\n"
            except Exception:
                pass
        return text


    elif ext in ('.xls', '.xlsx'):
        import tempfile
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
        try:
            tmp.write(content)
            tmp.close()
            df = pd.read_excel(tmp.name, engine='openpyxl' if ext=='.xlsx' else None)
            return df.head(50).to_csv(index=False)
        except Exception as e:
            logging.warning(f"Excel parse error: {e}")
            return ""
        finally:
            try: os.unlink(tmp.name)
            except Exception: pass

    elif ext == '.pptx':
        import tempfile
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
        try:
            tmp.write(content)
            tmp.close()
            prs = Presentation(tmp.name)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)
            return "\n".join(lines)
        except Exception as e:
            logging.warning(f"PPTX parse error: {e}")
            return ""
        finally:
            try: os.unlink(tmp.name)
            except Exception: pass

    elif ext == '.rtf':
        try:
            return rtf_to_text(content.decode(errors="ignore"))
        except Exception as e:
            logging.warning(f"RTF parse error: {e}")
            return ""

    elif ext in ('.odt', '.ods'):
        import tempfile
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
        try:
            tmp.write(content)
            tmp.close()
            doc = odf_load(tmp.name)
            paragraphs = []
            for elem in doc.getElementsByType(P):
                paragraphs.append(str(elem))
            return "\n".join(paragraphs)
        except Exception as e:
            logging.warning(f"ODF parse error: {e}")
            return ""
        finally:
            try: os.unlink(tmp.name)
            except Exception: pass

    elif ext == '.eml':
        try:
            msg = email.message_from_bytes(content, policy=email.policy.default)
            text_parts = []
            if msg['Subject']:
                text_parts.append("Subject: " + str(msg['Subject']))
            if msg['From']:
                text_parts.append("From: " + str(msg['From']))
            if msg['To']:
                text_parts.append("To: " + str(msg['To']))
            if msg['Date']:
                text_parts.append("Date: " + str(msg['Date']))
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body += str(part.get_content())
            else:
                body = msg.get_content()
            text_parts.append(str(body))
            return "\n".join(text_parts)
        except Exception as e:
            logging.warning(f"EML parse error: {e}")
            return ""

    elif ext == '.msg' and extract_msg is not None:
        import tempfile
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
        try:
            tmp.write(content)
            tmp.close()
            msg = extract_msg.Message(tmp.name)
            text = f"Subject: {msg.subject}\nFrom: {msg.sender}\nTo: {msg.to}\nDate: {msg.date}\n\n{msg.body}"
            return text
        except Exception as e:
            logging.warning(f"MSG parse error: {e}")
            return ""
        finally:
            try: os.unlink(tmp.name)
            except Exception: pass

    return "[Unknown/unsupported filetype]"

# def build_gpt_system_prompt(COMPLIANCE_MATRIX):
#     instructions = (
#         "You are an expert compliance auditor. Using the standards and fields below, "
#         "analyze each document for regulated data.\n"
#         "For each applicable compliance standard, list:\n"
#         "* The standard\n* The jurisdiction\n* Which relevant data types occur in this file.\n"
#         "If regulated data is present, please assess and provide risk rating of High, Medium or Low 'Risk rating is:'Risk\n"
#         "Based on your assessment at individual compliance standard, please provide overall risk rating. Overall rating should be the highest rating of risk rated across all compliance standards\n"
#     )
#     for entry in COMPLIANCE_MATRIX:
#         instructions += f"Standard: {entry['standard']} | Jurisdiction: {entry['jurisdiction']} | Fields: {', '.join(entry['fields'])}\n"
#     return instructions


def build_gpt_system_prompt(COMPLIANCE_MATRIX):
    instructions = (
        "You are an expert compliance auditor and cyber expert. Based on your expert knowledge and using the table below for standards and regulated data fields as a reference,\n"
        "analyze the provided document (or data) for regulated data types and regulated data fields. When reviewing, pay close attention to the following and document your reasoning:\n"
        "- If the file contains or appears to contain only training material, instructional content with fake/sample/demo data, or random values (not real production data), clearly indicate this in your results and reasoning.\n"
        "- If regulated fields are present but there are no real (production or sensitive) values associated, rate the risk and data classification LOWER than if real production data is present.\n"
        "- If the data is explicitly marked as sample data, test records, or is clearly intended for demonstration or training, reduce the file's risk and sensitivity rating accordingly and document this in your explanation.\n"
        "- Only assign High or Medium risk and 'sensitive' classifications where there is strong evidence of actual regulated data populated in those fields.\n"
        "\n"
        "For every compliance standard with a match, create a results array entry containing:\n"
        "  - \"standard\": name of the standard\n"
        "  - \"jurisdiction\": the law's jurisdiction/region\n"
        "  - \"detected_fields\": list of regulated field names found in this document\n"
        "  - \"risk\": High, Medium, Low, or None, based on the actual data detected and the nature of that data (with sample/demo/fake data as Low or None)\n"
        "After reviewing all, output the following as ordered fields at the END of the JSON object:\n"
        "  - \"auditor_agent_view\": In 2-3 sentences, explain your reasoning (plain English, for auditors), specifically highlight if the content appears to be training/sample/demo, or fake/random data.\n"
        "  - \"likely_data_subject_area\": State data subject area, e.g., customer data, financial records, health information, etc.\n"
        "  - \"data_classification\": Classify Data, e.g., Highly Sensitive, Confidential, Internal Use, Public etc. Use a lower classification if the file only contains training/sample/demo data or regulated fields without actual sensitive values.\n"
        "  - \"overall_risk_rating\": string, High, Medium, Low, or None (the highest risk found overall; use Low or None if only sample/demo/empty data)\n"
        "  - \"hacker_interest\": State How Hacker or bad actors will potentially exploit this information if they get access to it. If the file only has demo/sample data, state 'Little to no value to hackers.'\n"
        "  - \"cyber_proposed_controls\": List proposed cyber controls leveraging ISO standard and NIST, based on data classification and risk rating as a cyber expert\n"
        "  - \"auditor_proposed_action\": string, short action phrase for remediation based on auditor expert knowledge, e.g., \"notify owner\", \"quarantine file\", \"no action required\" etc. For training/sample/demo data, recommend \"no action required\" or similar.\n"
        "Respond with a SINGLE pure JSON object only, in this format:\n"
        "{\n"
        "  \"results\": [\n"
        "     {\"standard\":\"\",\"jurisdiction\":\"\",\"detected_fields\":[],\"risk\":\"\",\"auditor_agent_view\":\"\"}, ...\n"
        "  ],\n"
        "  \"auditor_agent_view\": \"\",\n"
        "  \"likely_data_subject_area\": \"\",\n"
        "  \"data_classification\": \"\",\n"
        "  \"overall_risk_rating\": \"\",\n"
        "  \"hacker_interest\": \"\",\n"
        "  \"cyber_proposed_controls\":: \"\",\n"
        "  \"auditor_proposed_action\": \"\"\n"
        "}\n"
        "Here is your standards/fields matrix:\n"
    )


    for entry in COMPLIANCE_MATRIX:
        instructions += (
            f"Standard: {entry['standard']} | Jurisdiction: {entry['jurisdiction']} | "
            f"Fields: {', '.join(entry['fields'])}\n"
        )
    return instructions


def gpt_classify_file(text, compliance_matrix):
    system_prompt = build_gpt_system_prompt(compliance_matrix)
    user_prompt = f"Classify the following file content for regulated data fields (see system context):\n\n{text[:CHUNK_SIZE]}"
    try:
        response = openai.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.0,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        logging.error(f"GPT API error: {e}")
        return f"GPT API error: {e}"

def iter_file_jsons(config_id, config_dir):
    # Pick up all *_m365_files_list_<config_id>.json in config_dir (user and SharePoint drive files)
    import re
    #pat = re.compile(r"([A-Za-z0-9._-]+)_m365_files_list_" + re.escape(config_id) + r"\.json$")
    pat = re.compile(r"(.+)_m365_files_list_" + re.escape(config_id) + r"\.json$")
    for file in os.listdir(config_dir):
        if pat.match(file):
            yield os.path.join(config_dir, file)

# ================
# MAIN LOGIC
# ================



def main():
    # --- Argument/config directory discovery ---
    if len(sys.argv) >= 2:
        config_id = sys.argv[1]
    else:
        config_id = input("Enter config_id: ").strip()
    search_root = sys.argv[2] if len(sys.argv) > 2 else '/home/cybersecai/htdocs/www.cybersecai.io/webhook/M365'

    logging.info(f"Startup: Looking for config folder '{config_id}' under root '{search_root}'")
    config_dir = find_config_folder(config_id, search_root)
    secrets = read_secrets(config_dir, config_id)
    compliance_matrix = load_compliance_matrix(config_dir)

    graph_dir = os.path.join(config_dir, "graph")
    os.makedirs(graph_dir, exist_ok=True)
    logging.info(f"Will write PII report(s) to: {graph_dir}")

    import openai
    openai.api_key = os.environ.get("OPENAI_API_KEY")
    if not openai.api_key:
        abort("Please set your OpenAI API key in the OPENAI_API_KEY environment variable.")

    token = get_access_token(secrets['TENANT_ID'], secrets['CLIENT_ID'], secrets['CLIENT_SECRET'])

    for json_file in iter_file_jsons(config_id, config_dir):
        logging.info(f"Processing file list: {json_file}")

        # Robust skip: empty/whitespace/bad JSON, or valid but empty list
        try:
            if os.path.getsize(json_file) == 0:
                logging.warning(f"File {json_file} is empty (0 bytes), skipping.")
                continue

            with open(json_file, "r", encoding="utf-8") as jf:
                content = jf.read()
            if not content or not content.strip():
                logging.warning(f"File {json_file} contains only whitespace, skipping.")
                continue

            try:
                file_records = json.loads(content)
            except json.JSONDecodeError as jde:
                logging.error(f"JSON decode error for {json_file}: {jde} (skipping)")
                continue

            if not file_records or not isinstance(file_records, list):
                logging.warning(f"File {json_file} contains no file records or is [] (empty list). Skipping.")
                continue

        except Exception as e:
            logging.error(f"Could not read {json_file}: {e}")
            continue

        output_list = []
        count_total = 0
        n_files, n_folders = 0, 0

        for rec in file_records:
            ext = os.path.splitext(rec.get('file_name', ''))[1].lower()
            # Only count as folder if it's a folder or not ALLOWED_EXT
            if ext not in ALLOWED_EXT or rec.get('file_type','').lower() == 'folder':
                n_folders += 1
                continue
            n_files += 1
            count_total += 1
            text = fetch_file_content(rec, token)
            text_cut = text[:MAX_OUT_CHARS] if text else ""
            if not text_cut.strip():
                logging.warning(f"No usable text for file {rec.get('file_name')}. Skipping.")
                continue
            logging.info(f"Classifying file: {rec.get('file_name')}")
            llm_response = gpt_classify_file(text_cut, compliance_matrix)
            record_with_llm = dict(rec)
            record_with_llm['llm_response'] = llm_response
            output_list.append(record_with_llm)

        logging.info(f"From {json_file}: {n_files} files to classify, {n_folders} folders/unsupported skipped.")
        if output_list:
            orig_basename = os.path.basename(json_file)
            output_fname = f"output_{orig_basename}"
            output_path = os.path.join(graph_dir, output_fname)
            try:
                with open(output_path, "w", encoding='utf-8') as outf:
                    json.dump(output_list, outf, indent=2)
                logging.info(f"*** Wrote PII report: {output_path} ({count_total} files checked)")
            except Exception as e:
                logging.error(f"Could not write output file: {output_path}. Error: {e}")
        else:
            logging.info(f"No PII found in any files for {json_file} (or no supported files present)")

if __name__ == '__main__':
    main()